"""Executor — consumes validated slide_spec JSON, emits a python-pptx Slide.

The executor is intentionally dumb: it dispatches each ShapeCall to the
registered shape-kit function with the spec's args. The Validator is what
guarantees the spec is internally consistent before we get here.
"""

from __future__ import annotations

import structlog
from pptx import Presentation
from pptx.util import Emu

from app.preso_pro.planning.slide_spec import DeckContext, SlideSpec
from app.preso_pro.shape_kit import get_function, is_registered

logger = structlog.get_logger()


def execute_slide(prs: Presentation, ctx: DeckContext, spec: SlideSpec):
    """Append a slide to `prs` rendered from `spec`."""
    blank_layout = prs.slide_layouts[6]  # 6 is the blank layout in default master
    slide = prs.slides.add_slide(blank_layout)

    # Background first so subsequent elements layer above it
    if spec.background:
        fn = get_function(spec.background.fn)
        if fn is None:
            logger.warn("preso_pro_unknown_bg_fn", fn=spec.background.fn)
        else:
            fn(slide, ctx, **spec.background.args)

    for elem in spec.elements:
        fn = get_function(elem.fn)
        if fn is None:
            logger.warn("preso_pro_unknown_fn", fn=elem.fn, slide=spec.slide_index)
            continue
        try:
            fn(slide, ctx, **elem.args)
        except TypeError as e:
            logger.error(
                "preso_pro_arg_mismatch",
                fn=elem.fn,
                args=elem.args,
                err=str(e),
                slide=spec.slide_index,
            )
        except Exception as e:
            logger.error(
                "preso_pro_render_failed",
                fn=elem.fn,
                err=str(e),
                slide=spec.slide_index,
            )

    return slide


def new_presentation() -> Presentation:
    """Create an empty 16:9 presentation."""
    prs = Presentation()
    prs.slide_width = Emu(12192000)
    prs.slide_height = Emu(6858000)
    return prs
