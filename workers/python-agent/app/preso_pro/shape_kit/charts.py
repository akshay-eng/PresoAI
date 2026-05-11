"""Native PPTX chart helpers — bar, line, pie, donut, area, scatter.

All charts are real native PowerPoint charts (data-bound), so the user can
edit them in Office. Colors come from the deck's locked palette.
"""

from __future__ import annotations

from typing import Any

from lxml import etree
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.slide import Slide
from pptx.util import Pt

from app.preso_pro.planning.slide_spec import DeckContext
from app.preso_pro.shape_kit._util import font_for_kind, hex_for_role, hex_to_rgb
from app.preso_pro.shape_kit.anchors import (
    SLIDE_HEIGHT_EMU,
    SLIDE_WIDTH_EMU,
    anchor_to_emu,
)

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
C_NS = "http://schemas.openxmlformats.org/drawingml/2006/chart"


def _chart_box(anchor: str, width_frac: float, height_frac: float) -> tuple[int, int, int, int]:
    x, y = anchor_to_emu(anchor)
    w = int(SLIDE_WIDTH_EMU * width_frac)
    h = int(SLIDE_HEIGHT_EMU * height_frac)
    safe = int(SLIDE_WIDTH_EMU * 0.05)
    left = max(safe, min(x - w // 2, SLIDE_WIDTH_EMU - w - safe))
    top = max(safe, min(y - h // 2, SLIDE_HEIGHT_EMU - h - safe))
    return left, top, w, h


def _palette_chart_colors(ctx: DeckContext, n: int) -> list[str]:
    """Return n hex colors from the palette accents — primary/accent_1/accent_2 cycled."""
    pool = [
        hex_for_role(ctx, "primary"),
        hex_for_role(ctx, "accent_1"),
        hex_for_role(ctx, "accent_2"),
        hex_for_role(ctx, "text_muted"),
    ]
    return [pool[i % len(pool)] for i in range(n)]


def _color_chart_series(chart: Any, colors_hex: list[str]) -> None:
    """Apply explicit fill colors to each series of a category chart."""
    for i, ser in enumerate(chart.series):
        hex_color = colors_hex[i % len(colors_hex)]
        r, g, b = hex_to_rgb(hex_color)
        fill = ser.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        ser.format.line.fill.background()


def _color_pie_points(chart: Any, colors_hex: list[str]) -> None:
    """Pie/donut: color each data point in series 0 separately."""
    if not chart.plots or not chart.plots[0].series:
        return
    series = chart.plots[0].series[0]
    points = list(series.points)
    for i, pt in enumerate(points):
        hex_color = colors_hex[i % len(colors_hex)]
        r, g, b = hex_to_rgb(hex_color)
        fill = pt.format.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(r, g, b)
        pt.format.line.fill.background()


def _style_chart_text(chart: Any, ctx: DeckContext, font_size_pt: int = 11) -> None:
    """Apply heading font + muted text role to chart text where exposed."""
    font_name = font_for_kind(ctx, "body")
    text_color = hex_for_role(ctx, "text_muted")
    r, g, b = hex_to_rgb(text_color)

    # Top-level chart text properties
    try:
        tx_pr = chart.text_frame.paragraphs[0].runs
        for run in tx_pr:
            run.font.name = font_name
            run.font.size = Pt(font_size_pt)
            run.font.color.rgb = RGBColor(r, g, b)
    except Exception:
        pass

    # Axis tick labels — best-effort via raw XML since python-pptx exposes
    # axes only on category charts.
    try:
        chart_part = chart._chartSpace
        for txPr in chart_part.iter(f"{{{C_NS}}}txPr"):
            # Replace defRPr / rPr font color + name for all visible text.
            for el in txPr.iter(f"{{{A_NS}}}defRPr"):
                el.set("sz", str(font_size_pt * 100))
                _ensure_solid_fill(el, r, g, b)
                _ensure_latin(el, font_name)
    except Exception:
        pass


def _ensure_solid_fill(rpr_el: Any, r: int, g: int, b: int) -> None:
    # remove any existing fill
    for tag in ("solidFill", "gradFill", "noFill"):
        existing = rpr_el.find(f"{{{A_NS}}}{tag}")
        if existing is not None:
            rpr_el.remove(existing)
    sf = etree.SubElement(rpr_el, f"{{{A_NS}}}solidFill")
    srgb = etree.SubElement(sf, f"{{{A_NS}}}srgbClr")
    srgb.set("val", f"{r:02X}{g:02X}{b:02X}")


def _ensure_latin(rpr_el: Any, font_name: str) -> None:
    existing = rpr_el.find(f"{{{A_NS}}}latin")
    if existing is not None:
        rpr_el.remove(existing)
    latin = etree.SubElement(rpr_el, f"{{{A_NS}}}latin")
    latin.set("typeface", font_name)


# ─── Chart helpers ────────────────────────────────────────────────────────

def bar_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None,
    series: list[dict] | None = None,
    orientation: str = "vertical",
    role_a: str = "primary", role_b: str = "accent_1",
) -> None:
    """A native PPTX bar chart (vertical or horizontal).

    series: list of {name, values} dicts. Example:
        [{"name": "2024", "values": [120, 140, 180]},
         {"name": "2025", "values": [150, 220, 260]}]
    """
    categories = categories or ["Q1", "Q2", "Q3", "Q4"]
    series = series or [{"name": "Revenue", "values": [10, 14, 18, 25]}]

    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(str(s.get("name", "Series")), tuple(s.get("values", [])))

    chart_type = (
        XL_CHART_TYPE.COLUMN_CLUSTERED if orientation == "vertical"
        else XL_CHART_TYPE.BAR_CLUSTERED
    )
    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(chart_type, left, top, w, h, cd)
    chart = graphic.chart

    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False

    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    colors = [hex_for_role(ctx, role_a), hex_for_role(ctx, role_b),
              hex_for_role(ctx, "accent_2"), hex_for_role(ctx, "text_muted")]
    _color_chart_series(chart, colors[: max(1, len(series))])
    _style_chart_text(chart, ctx)


def stacked_bar_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None,
    series: list[dict] | None = None,
) -> None:
    categories = categories or ["Q1", "Q2", "Q3", "Q4"]
    series = series or [
        {"name": "New", "values": [40, 55, 70, 85]},
        {"name": "Returning", "values": [60, 70, 80, 95]},
    ]
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(str(s.get("name", "Series")), tuple(s.get("values", [])))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    _color_chart_series(chart, _palette_chart_colors(ctx, len(series)))
    _style_chart_text(chart, ctx)


def line_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None,
    series: list[dict] | None = None,
) -> None:
    categories = categories or ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    series = series or [{"name": "MRR", "values": [10, 14, 18, 24, 32, 47]}]
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(str(s.get("name", "Series")), tuple(s.get("values", [])))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.LINE, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    colors = _palette_chart_colors(ctx, len(series))
    for i, ser in enumerate(chart.series):
        hex_color = colors[i % len(colors)]
        r, g, b = hex_to_rgb(hex_color)
        ser.format.line.color.rgb = RGBColor(r, g, b)
        ser.format.line.width = Pt(3)
    _style_chart_text(chart, ctx)


def area_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None,
    series: list[dict] | None = None,
) -> None:
    categories = categories or ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    series = series or [{"name": "Users", "values": [120, 180, 280, 420, 600, 880]}]
    cd = CategoryChartData()
    cd.categories = categories
    for s in series:
        cd.add_series(str(s.get("name", "Series")), tuple(s.get("values", [])))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.AREA, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = len(series) > 1

    _color_chart_series(chart, _palette_chart_colors(ctx, len(series)))
    _style_chart_text(chart, ctx)


def pie_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.40, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None, values: list[float] | None = None,
) -> None:
    categories = categories or ["Enterprise", "Mid-market", "SMB"]
    values = values or [55, 30, 15]
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Share", tuple(values))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.PIE, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    _color_pie_points(chart, _palette_chart_colors(ctx, len(values)))
    _style_chart_text(chart, ctx)


def scatter_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.55,
    title: str = "", series: list[dict] | None = None,
) -> None:
    """Native scatter (XY) chart. Useful for distribution / correlation.

    series: list of {name, points: [[x, y], ...]} dicts.
    """
    series = series or [{"name": "Customers",
                         "points": [[10, 12], [22, 18], [35, 30], [48, 42], [60, 55]]}]
    cd = XyChartData()
    for s in series:
        ser = cd.add_series(str(s.get("name", "Series")))
        for pt in s.get("points") or []:
            if isinstance(pt, (list, tuple)) and len(pt) >= 2:
                ser.add_data_point(float(pt[0]), float(pt[1]))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.XY_SCATTER, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = len(series) > 1

    colors = _palette_chart_colors(ctx, len(series))
    for i, ser in enumerate(chart.series):
        hex_color = colors[i % len(colors)]
        r, g, b = hex_to_rgb(hex_color)
        ser.format.line.color.rgb = RGBColor(r, g, b)
    _style_chart_text(chart, ctx)


def donut_chart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.40, height_frac: float = 0.55,
    title: str = "", categories: list[str] | None = None, values: list[float] | None = None,
) -> None:
    categories = categories or ["Enterprise", "Mid-market", "SMB"]
    values = values or [55, 30, 15]
    cd = CategoryChartData()
    cd.categories = categories
    cd.add_series("Share", tuple(values))

    left, top, w, h = _chart_box(anchor, width_frac, height_frac)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, left, top, w, h, cd)
    chart = graphic.chart
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.text = title
    else:
        chart.has_title = False
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False

    _color_pie_points(chart, _palette_chart_colors(ctx, len(values)))
    _style_chart_text(chart, ctx)
