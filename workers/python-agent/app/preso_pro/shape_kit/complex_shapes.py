"""Complex marketing shapes — arrows, stars, hexagons, chevrons, and the
process/flow composites that string them together.

These give the LLM richer vocabulary than rectangles + circles. Process
composites take a list of step strings and render numbered/arrow chains.
"""

from __future__ import annotations

import math
from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from app.preso_pro.planning.slide_spec import DeckContext
from app.preso_pro.shape_kit._util import (
    font_for_kind,
    hex_for_role,
    hex_to_rgb,
    size_for_tier,
)
from app.preso_pro.shape_kit.anchors import (
    SLIDE_HEIGHT_EMU,
    SLIDE_WIDTH_EMU,
    anchor_to_emu,
)
from app.preso_pro.shape_kit import ooxml_helpers as ox

EMU_PER_PT = 12700
SAFE_MARGIN_EMU = int(SLIDE_WIDTH_EMU * 0.05)


# ─── Single complex primitives ────────────────────────────────────────────

def right_arrow(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", length: int = 200, height: int = 60,
    role: str = "primary", opacity: float = 1.0, direction: str = "right",
) -> None:
    """A pointing arrow shape. direction: right|left|up|down."""
    x, y = anchor_to_emu(anchor)
    length_emu = int(length * EMU_PER_PT)
    height_emu = int(height * EMU_PER_PT)
    shape_map = {
        "right": MSO_SHAPE.RIGHT_ARROW,
        "left": MSO_SHAPE.LEFT_ARROW,
        "up": MSO_SHAPE.UP_ARROW,
        "down": MSO_SHAPE.DOWN_ARROW,
    }
    if direction in ("up", "down"):
        w, h = height_emu, length_emu
    else:
        w, h = length_emu, height_emu
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(shape_map.get(direction, MSO_SHAPE.RIGHT_ARROW),
                                   left, top, w, h)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def star(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "upper-right", size: int = 80,
    points: int = 5, role: str = "accent_1", opacity: float = 1.0,
) -> None:
    """A star shape — 5, 6, or 8 points."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    shape_map = {
        4: MSO_SHAPE.STAR_4_POINT,
        5: MSO_SHAPE.STAR_5_POINT,
        6: MSO_SHAPE.STAR_6_POINT,
        7: MSO_SHAPE.STAR_7_POINT,
        8: MSO_SHAPE.STAR_8_POINT,
        10: MSO_SHAPE.STAR_10_POINT,
    }
    kind = shape_map.get(points, MSO_SHAPE.STAR_5_POINT)
    shape = slide.shapes.add_shape(kind, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def hexagon(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 120,
    role: str = "primary", opacity: float = 1.0,
) -> None:
    """A regular hexagon."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.HEXAGON, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def chevron(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", length: int = 150, height: int = 60,
    role: str = "accent_1", opacity: float = 1.0,
) -> None:
    """A chevron arrow — used in process flows for connection."""
    x, y = anchor_to_emu(anchor)
    w = int(length * EMU_PER_PT)
    h = int(height * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, w, h)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def gear(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 100,
    role: str = "primary", opacity: float = 1.0,
) -> None:
    """A gear/cog icon shape."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.GEAR_9, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def heart(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 100,
    role: str = "accent_1", opacity: float = 1.0,
) -> None:
    """Heart icon — for love/affinity slides."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.HEART, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def shield(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 100,
    role: str = "primary", opacity: float = 1.0,
) -> None:
    """Pentagon as a shield icon — for security/trust slides."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.PENTAGON, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


def checkmark_badge(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 80,
    role_circle: str = "accent_2", role_check: str = "text_inverse",
) -> None:
    """A circle with a checkmark inside — affirmation/confirmed accent.

    Renders as an oval with a unicode checkmark text on top."""
    x, y = anchor_to_emu(anchor)
    side = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - side // 2, SLIDE_WIDTH_EMU - side - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - side // 2, SLIDE_HEIGHT_EMU - side - SAFE_MARGIN_EMU))
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, side, side)
    circle.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role_circle))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(r, g, b)
    # check glyph
    _add_text_centered(
        slide, ctx, text="✓",
        left=left, top=top + int(side * 0.10),
        width=side, height=int(side * 0.80),
        tier="h1", role=role_check, bold=True, font_kind="heading",
    )


def lightning_bolt(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", size: int = 120,
    role: str = "accent_1", opacity: float = 1.0,
) -> None:
    """Lightning bolt — emphasis / energy decorator."""
    x, y = anchor_to_emu(anchor)
    w = int(size * EMU_PER_PT * 0.6)
    h = int(size * EMU_PER_PT)
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.LIGHTNING_BOLT, left, top, w, h)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    if opacity < 1.0:
        ox.apply_alpha(shape, opacity)


# ─── Process / flow composites ────────────────────────────────────────────

def _add_text_centered(
    slide: Slide, ctx: DeckContext, *, text: str,
    left: int, top: int, width: int, height: int,
    tier: str = "body", role: str = "text_primary",
    bold: bool = False, font_kind: str = "body",
) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, font_kind)
    run.font.size = Pt(size_for_tier(ctx, tier))
    run.font.bold = bold
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    run.font.color.rgb = RGBColor(r, g, b)


_AVG_CHAR_W_EMU_PER_PT = 11000  # rough character width estimate
_EMU_PER_PT_LOCAL = 12700


def _add_text_fitted(
    slide: Slide, ctx: DeckContext, *, text: str,
    left: int, top: int, width: int, height: int,
    max_tier: str = "h2", min_pt: int = 10, role: str = "text_primary",
    bold: bool = False, font_kind: str = "body",
) -> None:
    """Like _add_text_centered, but shrinks the font to fit width AND height.

    Use this inside small shape interiors (cycle nodes, swot headers, radial hub)
    where rigid tier sizing causes wrap or overflow.
    """
    max_pt = size_for_tier(ctx, max_tier)
    chars = max(1, len(text))
    fit_w = max(min_pt, int(width / (chars * _AVG_CHAR_W_EMU_PER_PT)))
    fit_h = max(min_pt, int((height * 0.70) / _EMU_PER_PT_LOCAL))
    size_pt = max(min_pt, min(max_pt, fit_w, fit_h))

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, font_kind)
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    run.font.color.rgb = RGBColor(r, g, b)


def arrow_flow(
    slide: Slide, ctx: DeckContext,
    *, steps: list[str], anchor: str = "center",
    role_arrow: str = "primary", role_text: str = "text_primary",
) -> None:
    """A horizontal sequence of right-arrows with text on each.

    steps: list of short strings (3-5 items).
    """
    n = max(2, min(5, len(steps)))
    steps = steps[:n]
    x, y = anchor_to_emu(anchor)

    total_w = int(SLIDE_WIDTH_EMU * 0.84)
    arrow_w = total_w // n
    arrow_h = int(SLIDE_HEIGHT_EMU * 0.16)
    overlap = int(arrow_w * 0.06)

    start_left = (SLIDE_WIDTH_EMU - total_w) // 2
    top = max(SAFE_MARGIN_EMU, min(y - arrow_h // 2, SLIDE_HEIGHT_EMU - arrow_h - SAFE_MARGIN_EMU))

    base_color = hex_for_role(ctx, role_arrow)
    accent_color = hex_for_role(ctx, "accent_1")

    for i in range(n):
        step_left = start_left + i * (arrow_w - overlap)
        # Use pentagon for "tail" arrows and chevron for the inner segments
        # for a clean stepped look. Last step is right_arrow (clear endpoint).
        kind = MSO_SHAPE.PENTAGON if i == 0 else (
            MSO_SHAPE.CHEVRON if i < n - 1 else MSO_SHAPE.RIGHT_ARROW
        )
        shape = slide.shapes.add_shape(kind, step_left, top, arrow_w, arrow_h)
        shape.line.fill.background()
        # alternate colors slightly for visual rhythm
        hex_color = base_color if i % 2 == 0 else accent_color
        r, g, b = hex_to_rgb(hex_color)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        _polish_node(shape, ctx, role_arrow if i % 2 == 0 else "accent_1",
                     vertical=False, shadow=True)

        # Text inside the step
        text_pad = int(arrow_w * 0.10)
        _add_text_centered(
            slide, ctx, text=steps[i],
            left=step_left + text_pad, top=top, width=arrow_w - 2 * text_pad,
            height=arrow_h, tier="body", role=role_text, bold=True, font_kind="body",
        )


def numbered_steps(
    slide: Slide, ctx: DeckContext,
    *, items: list[dict], anchor: str = "center",
    role_circle: str = "primary", role_label: str = "text_primary",
) -> None:
    """Horizontal sequence of numbered circles with labels underneath.

    items: list of {label} dicts. Numbers are auto-generated 1..N.
    Connecting chevrons are drawn between circles.
    """
    n = max(2, min(6, len(items)))
    items = items[:n]
    x, y = anchor_to_emu(anchor)

    total_w = int(SLIDE_WIDTH_EMU * 0.84)
    cell_w = total_w // n
    circle_size = int(min(cell_w * 0.45, SLIDE_HEIGHT_EMU * 0.18))
    label_size_pt = size_for_tier(ctx, "body")
    label_h = int(label_size_pt * EMU_PER_PT * 1.5 * 2)
    chevron_w = int(cell_w * 0.18)
    chevron_h = int(circle_size * 0.30)

    start_left = (SLIDE_WIDTH_EMU - total_w) // 2
    block_h = circle_size + Emu(80000) + label_h
    top_circle = max(SAFE_MARGIN_EMU, min(y - block_h // 2,
                                           SLIDE_HEIGHT_EMU - block_h - SAFE_MARGIN_EMU))
    top_label = top_circle + circle_size + Emu(80000)

    cr, cg, cb = hex_to_rgb(hex_for_role(ctx, role_circle))
    accent_r, accent_g, accent_b = hex_to_rgb(hex_for_role(ctx, "accent_1"))
    text_inv_r, text_inv_g, text_inv_b = hex_to_rgb(hex_for_role(ctx, "text_inverse"))

    for i, item in enumerate(items):
        cell_left = start_left + i * cell_w
        cx = cell_left + cell_w // 2
        circle_left = cx - circle_size // 2

        # Circle — gradient + shadow for SmartArt-quality polish
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, circle_left, top_circle,
                                      circle_size, circle_size)
        oval.line.fill.background()
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(cr, cg, cb)
        _polish_node(oval, ctx, role_circle)

        # Number text inside circle
        _add_text_centered(
            slide, ctx, text=str(i + 1),
            left=circle_left, top=top_circle + int(circle_size * 0.18),
            width=circle_size, height=int(circle_size * 0.7),
            tier="h2", role="text_inverse", bold=True, font_kind="heading",
        )

        # Chevron connecting to next step — bigger + gradient
        if i < n - 1:
            bigger_chevron_w = int(cell_w * 0.22)
            bigger_chevron_h = int(circle_size * 0.42)
            ch_left = circle_left + circle_size + int(cell_w * 0.04)
            ch_top = top_circle + (circle_size - bigger_chevron_h) // 2
            ch = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, ch_left, ch_top,
                                        bigger_chevron_w, bigger_chevron_h)
            ch.line.fill.background()
            ch.fill.solid()
            ch.fill.fore_color.rgb = RGBColor(accent_r, accent_g, accent_b)
            _polish_node(ch, ctx, "accent_1", vertical=False, shadow=False)

        # Label below
        _add_text_centered(
            slide, ctx, text=str(item.get("label", "")),
            left=cell_left, top=top_label, width=cell_w, height=label_h,
            tier="body", role=role_label, bold=False, font_kind="body",
        )


def feature_grid(
    slide: Slide, ctx: DeckContext,
    *, items: list[dict], anchor: str = "center",
    role_card: str = "surface", role_icon: str = "primary",
    role_title: str = "text_primary", role_body: str = "text_muted",
) -> None:
    """Marketing-deck core-features layout: 3 or 4 cards, each with an icon
    glyph (rounded rectangle), a heading, and a one-line description.

    items: list of {icon, title, body} dicts. icon is a 1-3 char string
    (emoji or letter). 3-4 items recommended.
    """
    n = max(2, min(4, len(items)))
    items = items[:n]
    x, y = anchor_to_emu(anchor)

    total_w = int(SLIDE_WIDTH_EMU * 0.84)
    gap = int(SLIDE_WIDTH_EMU * 0.025)
    card_w = (total_w - gap * (n - 1)) // n
    card_h = int(SLIDE_HEIGHT_EMU * 0.42)
    icon_size = int(card_w * 0.22)
    title_size = size_for_tier(ctx, "h2")
    body_size = size_for_tier(ctx, "body")

    start_left = (SLIDE_WIDTH_EMU - total_w) // 2
    top = max(SAFE_MARGIN_EMU, min(y - card_h // 2,
                                    SLIDE_HEIGHT_EMU - card_h - SAFE_MARGIN_EMU))

    cr, cg, cb = hex_to_rgb(hex_for_role(ctx, role_card))
    icon_r, icon_g, icon_b = hex_to_rgb(hex_for_role(ctx, role_icon))
    title_r, title_g, title_b = hex_to_rgb(hex_for_role(ctx, role_title))
    body_r, body_g, body_b = hex_to_rgb(hex_for_role(ctx, role_body))
    text_inv_hex = hex_for_role(ctx, "text_inverse")
    text_inv_r, text_inv_g, text_inv_b = hex_to_rgb(text_inv_hex)

    inner_pad = int(card_w * 0.10)

    for i, item in enumerate(items):
        cell_left = start_left + i * (card_w + gap)
        # Card backdrop — subtle gradient + soft shadow
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cell_left, top, card_w, card_h,
        )
        card.line.fill.background()
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(cr, cg, cb)
        _polish_node(card, ctx, role_card, light_pct=0.06, dark_pct=-0.04)

        # Icon container — rounded rectangle in primary, gradient + shadow
        icon_left = cell_left + inner_pad
        icon_top = top + inner_pad
        icon_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            icon_left, icon_top, icon_size, icon_size,
        )
        icon_box.line.fill.background()
        icon_box.fill.solid()
        icon_box.fill.fore_color.rgb = RGBColor(icon_r, icon_g, icon_b)
        _polish_node(icon_box, ctx, role_icon, shadow=False)

        # Icon glyph (text inside the icon container)
        glyph = str(item.get("icon", "·"))[:3]
        _add_text_centered(
            slide, ctx, text=glyph,
            left=icon_left, top=icon_top + int(icon_size * 0.10),
            width=icon_size, height=int(icon_size * 0.85),
            tier="h2", role="text_inverse", bold=True, font_kind="heading",
        )

        # Title (below the icon)
        title_top = icon_top + icon_size + int(SLIDE_HEIGHT_EMU * 0.025)
        title_h = int(title_size * EMU_PER_PT * 1.4 * 2)
        tbox = slide.shapes.add_textbox(
            cell_left + inner_pad, title_top,
            card_w - 2 * inner_pad, title_h,
        )
        ttf = tbox.text_frame
        ttf.word_wrap = True
        tp = ttf.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        tr.text = str(item.get("title", ""))
        tr.font.name = font_for_kind(ctx, "heading")
        tr.font.size = Pt(title_size)
        tr.font.bold = True
        tr.font.color.rgb = RGBColor(title_r, title_g, title_b)

        # Body (below the title)
        body_top = title_top + title_h + int(SLIDE_HEIGHT_EMU * 0.010)
        body_h = card_h - (body_top - top) - inner_pad
        if body_h > int(body_size * EMU_PER_PT):
            bbox = slide.shapes.add_textbox(
                cell_left + inner_pad, body_top,
                card_w - 2 * inner_pad, body_h,
            )
            btf = bbox.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.LEFT
            br = bp.add_run()
            br.text = str(item.get("body", ""))
            br.font.name = font_for_kind(ctx, "body")
            br.font.size = Pt(body_size)
            br.font.color.rgb = RGBColor(body_r, body_g, body_b)


def _connect_line(
    slide: Slide, ctx: DeckContext,
    *, x1: int, y1: int, x2: int, y2: int,
    role: str = "text_muted", width_pt: float = 1.5,
) -> Any:
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    line.line.color.rgb = RGBColor(r, g, b)
    line.line.width = Pt(width_pt)
    return line


# ─── Polish helpers — what makes shapes look like real PowerPoint SmartArt ──

def _shift_hex(hex_str: str, percent: float) -> str:
    """Lighten (positive %) or darken (negative %) a hex color."""
    r, g, b = hex_to_rgb(hex_str)
    if percent > 0:
        r = int(r + (255 - r) * percent)
        g = int(g + (255 - g) * percent)
        b = int(b + (255 - b) * percent)
    else:
        r = int(r * (1 + percent))
        g = int(g * (1 + percent))
        b = int(b * (1 + percent))
    return "#{:02X}{:02X}{:02X}".format(
        max(0, min(255, r)), max(0, min(255, g)), max(0, min(255, b))
    )


def _polish_node(
    shape: Any, ctx: DeckContext, role: str,
    *, vertical: bool = True, shadow: bool = True,
    light_pct: float = 0.18, dark_pct: float = -0.12,
) -> None:
    """Upgrade a flat shape to gradient + shadow — the 'pro' look.

    Replaces the existing solid fill with a 2-stop linear gradient (lighter
    top → darker bottom by default) and adds a subtle outer drop shadow.
    """
    base = hex_for_role(ctx, role)
    light = _shift_hex(base, light_pct)
    dark = _shift_hex(base, dark_pct)
    angle = 90 if vertical else 0
    try:
        ox.apply_linear_gradient(shape, light, dark, angle_deg=angle)
    except Exception:
        pass
    if shadow:
        try:
            ox.apply_outer_shadow(
                shape,
                blur=63500,        # ~5pt blur — soft, not muddy
                distance=25400,    # ~2pt offset
                direction=5400000, # 90° = straight down
                hex_color="#000000",
                alpha=35000,       # 35% — present but not heavy
            )
        except Exception:
            pass


def _polish_strip(
    shape: Any, ctx: DeckContext, role_a: str, role_b: str,
    *, angle_deg: int = 0, shadow: bool = False,
) -> None:
    """Two-color gradient strip (used for headers, spines, accent bars)."""
    a = hex_for_role(ctx, role_a)
    b = hex_for_role(ctx, role_b)
    try:
        ox.apply_linear_gradient(shape, a, b, angle_deg=angle_deg)
    except Exception:
        pass
    if shadow:
        try:
            ox.apply_outer_shadow(shape, blur=50800, distance=12700, alpha=30000)
        except Exception:
            pass


def cycle_diagram(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center",
    role_node: str = "primary", role_arrow: str = "accent_1",
    role_text: str = "text_inverse",
) -> None:
    """N-step cyclical flow — nodes around a circle with rotated arrows between
    each pair. items: 3-6 short strings."""
    n = max(3, min(6, len(items)))
    items = items[:n]
    cx, cy = anchor_to_emu(anchor)

    radius = int(min(SLIDE_WIDTH_EMU * 0.22, SLIDE_HEIGHT_EMU * 0.32))
    node_size = int(SLIDE_HEIGHT_EMU * 0.20)
    arrow_len = int(min(SLIDE_WIDTH_EMU * 0.10, radius * 0.65))
    arrow_h = int(arrow_len * 0.50)

    nr, ng, nb = hex_to_rgb(hex_for_role(ctx, role_node))
    ar, ag, ab = hex_to_rgb(hex_for_role(ctx, role_arrow))

    # Place nodes
    node_centers: list[tuple[int, int]] = []
    for i in range(n):
        theta = 2 * math.pi * i / n - math.pi / 2  # screen-space angle
        nx = cx + int(radius * math.cos(theta))
        ny = cy + int(radius * math.sin(theta))
        node_centers.append((nx, ny))

        oval = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            nx - node_size // 2, ny - node_size // 2, node_size, node_size,
        )
        oval.line.fill.background()
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(nr, ng, nb)
        _polish_node(oval, ctx, role_node)

        _add_text_fitted(
            slide, ctx, text=items[i],
            left=nx - node_size // 2 + Emu(60000),
            top=ny - int(node_size * 0.35),
            width=node_size - Emu(120000), height=int(node_size * 0.7),
            max_tier="body", min_pt=10, role=role_text, bold=True, font_kind="heading",
        )

    # Place arrows at midpoints between adjacent nodes (clockwise)
    for i in range(n):
        theta_m = 2 * math.pi * (i + 0.5) / n - math.pi / 2
        # Pull the arrow slightly outside the node radius for a clean look
        ar_radius = int(radius * 0.95)
        ax = cx + int(ar_radius * math.cos(theta_m))
        ay = cy + int(ar_radius * math.sin(theta_m))
        rotation_deg = math.degrees(theta_m + math.pi / 2)

        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            ax - arrow_len // 2, ay - arrow_h // 2, arrow_len, arrow_h,
        )
        arrow.rotation = rotation_deg
        arrow.line.fill.background()
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(ar, ag, ab)
        _polish_node(arrow, ctx, role_arrow, vertical=False, shadow=False)


def org_chart(
    slide: Slide, ctx: DeckContext,
    *, root: str, children: list[str], anchor: str = "center",
    role_root: str = "primary", role_child: str = "surface",
    role_root_text: str = "text_inverse", role_child_text: str = "text_primary",
    role_line: str = "text_muted",
) -> None:
    """Hierarchy: one root node with N children boxes. children: 2-5 strings."""
    n = max(2, min(5, len(children)))
    children = children[:n]
    cx, cy = anchor_to_emu(anchor)

    block_w = int(SLIDE_WIDTH_EMU * 0.84)
    box_h = int(SLIDE_HEIGHT_EMU * 0.16)
    gap_x = int(SLIDE_WIDTH_EMU * 0.02)
    child_w = (block_w - gap_x * (n - 1)) // n
    root_w = max(int(block_w * 0.32), int(child_w * 1.1))
    vertical_gap = int(SLIDE_HEIGHT_EMU * 0.14)

    block_h = box_h + vertical_gap + box_h
    top_root = max(SAFE_MARGIN_EMU,
                   min(cy - block_h // 2, SLIDE_HEIGHT_EMU - block_h - SAFE_MARGIN_EMU))
    top_children = top_root + box_h + vertical_gap

    root_left = cx - root_w // 2
    start_left = cx - block_w // 2

    # Root box — gradient + shadow
    rr, rg, rb = hex_to_rgb(hex_for_role(ctx, role_root))
    root_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, root_left, top_root, root_w, box_h,
    )
    root_box.line.fill.background()
    root_box.fill.solid()
    root_box.fill.fore_color.rgb = RGBColor(rr, rg, rb)
    _polish_node(root_box, ctx, role_root)
    _add_text_centered(
        slide, ctx, text=root,
        left=root_left + Emu(80000), top=top_root + int(box_h * 0.20),
        width=root_w - Emu(160000), height=int(box_h * 0.65),
        tier="h2", role=role_root_text, bold=True, font_kind="heading",
    )

    # Children + connecting lines (lines first so boxes draw on top)
    cr, cg, cb = hex_to_rgb(hex_for_role(ctx, role_child))
    root_bottom_x = cx
    root_bottom_y = top_root + box_h

    for i in range(n):
        cell_left = start_left + i * (child_w + gap_x)
        cell_cx = cell_left + child_w // 2

        # Connector first
        _connect_line(
            slide, ctx,
            x1=root_bottom_x, y1=root_bottom_y,
            x2=cell_cx, y2=top_children,
            role="primary", width_pt=2.0,
        )

        child_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cell_left, top_children, child_w, box_h,
        )
        child_box.line.fill.background()
        child_box.fill.solid()
        child_box.fill.fore_color.rgb = RGBColor(cr, cg, cb)
        _polish_node(child_box, ctx, role_child)
        _add_text_centered(
            slide, ctx, text=children[i],
            left=cell_left + Emu(60000), top=top_children + int(box_h * 0.22),
            width=child_w - Emu(120000), height=int(box_h * 0.6),
            tier="body", role=role_child_text, bold=True, font_kind="heading",
        )


def timeline(
    slide: Slide, ctx: DeckContext,
    *, items: list[dict], anchor: str = "center",
    role_line: str = "primary", role_marker: str = "accent_1",
    role_text: str = "text_primary", role_date: str = "text_muted",
) -> None:
    """Horizontal timeline with date markers. items: list of {date, label} dicts (3-6)."""
    n = max(3, min(6, len(items)))
    items = items[:n]
    cx, cy = anchor_to_emu(anchor)

    line_w = int(SLIDE_WIDTH_EMU * 0.84)
    line_left = cx - line_w // 2
    line_thickness = Emu(50000)
    marker_size = int(SLIDE_HEIGHT_EMU * 0.06)
    line_y = cy

    label_h = int(SLIDE_HEIGHT_EMU * 0.10)
    date_h = int(SLIDE_HEIGHT_EMU * 0.05)

    # Spine — gradient bar (primary → accent_1) for rhythm
    spine = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, line_left, line_y - line_thickness // 2,
        line_w, line_thickness,
    )
    spine.line.fill.background()
    sr, sg, sb = hex_to_rgb(hex_for_role(ctx, role_line))
    spine.fill.solid()
    spine.fill.fore_color.rgb = RGBColor(sr, sg, sb)
    _polish_strip(spine, ctx, role_line, "accent_1", angle_deg=0)

    mr, mg, mb = hex_to_rgb(hex_for_role(ctx, role_marker))

    spacing = line_w // (n - 1) if n > 1 else 0
    for i, item in enumerate(items):
        mx = line_left + i * spacing
        my = line_y

        # Marker (circle) — gradient + shadow
        slide_marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            mx - marker_size // 2, my - marker_size // 2,
            marker_size, marker_size,
        )
        slide_marker.line.fill.background()
        slide_marker.fill.solid()
        slide_marker.fill.fore_color.rgb = RGBColor(mr, mg, mb)
        _polish_node(slide_marker, ctx, role_marker)

        # Alternate above/below: even indices above, odd below
        text_w = int(spacing * 0.9) if spacing > 0 else int(SLIDE_WIDTH_EMU * 0.20)
        text_left = mx - text_w // 2

        if i % 2 == 0:
            # Date above marker, label above date
            date_top = my - marker_size - date_h - Emu(30000)
            label_top = date_top - label_h - Emu(20000)
        else:
            date_top = my + marker_size + Emu(30000)
            label_top = date_top + date_h + Emu(20000)

        _add_text_centered(
            slide, ctx, text=str(item.get("date", "")),
            left=text_left, top=date_top, width=text_w, height=date_h,
            tier="caption", role=role_date, bold=True, font_kind="heading",
        )
        _add_text_centered(
            slide, ctx, text=str(item.get("label", "")),
            left=text_left, top=label_top, width=text_w, height=label_h,
            tier="body", role=role_text, bold=False, font_kind="body",
        )


def vertical_process(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center",
    role_box: str = "primary", role_arrow: str = "accent_1",
    role_text: str = "text_inverse",
) -> None:
    """Vertically stacked steps with down-chevrons between them. items: 2-5 strings."""
    n = max(2, min(5, len(items)))
    items = items[:n]
    cx, cy = anchor_to_emu(anchor)

    box_w = int(SLIDE_WIDTH_EMU * 0.45)
    arrow_h = int(SLIDE_HEIGHT_EMU * 0.05)
    arrow_w = int(box_w * 0.18)
    total_h = int(SLIDE_HEIGHT_EMU * 0.78)
    box_h = (total_h - arrow_h * (n - 1)) // n

    block_top = max(SAFE_MARGIN_EMU,
                    min(cy - total_h // 2, SLIDE_HEIGHT_EMU - total_h - SAFE_MARGIN_EMU))
    box_left = cx - box_w // 2

    br, bg, bb = hex_to_rgb(hex_for_role(ctx, role_box))
    ar, ag, ab = hex_to_rgb(hex_for_role(ctx, role_arrow))

    for i in range(n):
        top = block_top + i * (box_h + arrow_h)
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, box_left, top, box_w, box_h,
        )
        box.line.fill.background()
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(br, bg, bb)
        _polish_node(box, ctx, role_box)
        _add_text_fitted(
            slide, ctx, text=items[i],
            left=box_left + Emu(80000), top=top + int(box_h * 0.18),
            width=box_w - Emu(160000), height=int(box_h * 0.7),
            max_tier="h2", min_pt=14, role=role_text, bold=True, font_kind="heading",
        )

        if i < n - 1:
            arrow_top = top + box_h
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                cx - arrow_w // 2, arrow_top, arrow_w, arrow_h,
            )
            arrow.line.fill.background()
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor(ar, ag, ab)
            _polish_node(arrow, ctx, role_arrow, shadow=False)


def radial_list(
    slide: Slide, ctx: DeckContext,
    *, hub: str, items: list[str], anchor: str = "center",
    role_hub: str = "primary", role_node: str = "surface",
    role_hub_text: str = "text_inverse", role_node_text: str = "text_primary",
    role_line: str = "text_muted",
) -> None:
    """Center hub with N radiating items connected by lines. items: 3-6 strings."""
    n = max(3, min(6, len(items)))
    items = items[:n]
    cx, cy = anchor_to_emu(anchor)

    radius = int(min(SLIDE_WIDTH_EMU * 0.26, SLIDE_HEIGHT_EMU * 0.36))
    hub_size = int(SLIDE_HEIGHT_EMU * 0.22)
    node_w = int(SLIDE_WIDTH_EMU * 0.20)
    node_h = int(SLIDE_HEIGHT_EMU * 0.10)

    hr, hg, hb = hex_to_rgb(hex_for_role(ctx, role_hub))
    nr, ng, nb = hex_to_rgb(hex_for_role(ctx, role_node))

    # Radiating nodes + connecting lines (lines first so nodes draw on top)
    for i in range(n):
        theta = 2 * math.pi * i / n - math.pi / 2
        nx = cx + int(radius * math.cos(theta))
        ny = cy + int(radius * math.sin(theta))

        # Thicker line in primary color — much more substantial than 1.25pt muted
        _connect_line(
            slide, ctx,
            x1=cx, y1=cy, x2=nx, y2=ny,
            role="primary", width_pt=2.5,
        )

        # Node pill — gradient + shadow
        nd = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            nx - node_w // 2, ny - node_h // 2, node_w, node_h,
        )
        nd.line.fill.background()
        nd.fill.solid()
        nd.fill.fore_color.rgb = RGBColor(nr, ng, nb)
        _polish_node(nd, ctx, role_node)
        _add_text_fitted(
            slide, ctx, text=items[i],
            left=nx - node_w // 2 + Emu(50000),
            top=ny - int(node_h * 0.35),
            width=node_w - Emu(100000), height=int(node_h * 0.7),
            max_tier="body", min_pt=10, role=role_node_text, bold=True, font_kind="heading",
        )

    # Hub on top so lines tuck behind it — gradient + shadow
    hub_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        cx - hub_size // 2, cy - hub_size // 2, hub_size, hub_size,
    )
    hub_shape.line.fill.background()
    hub_shape.fill.solid()
    hub_shape.fill.fore_color.rgb = RGBColor(hr, hg, hb)
    _polish_node(hub_shape, ctx, role_hub)
    _add_text_fitted(
        slide, ctx, text=hub,
        left=cx - hub_size // 2 + Emu(60000),
        top=cy - int(hub_size * 0.35),
        width=hub_size - Emu(120000), height=int(hub_size * 0.7),
        max_tier="h2", min_pt=14, role=role_hub_text, bold=True, font_kind="heading",
    )


def swot_quad(
    slide: Slide, ctx: DeckContext,
    *, quadrants: list[dict], anchor: str = "center",
    role_a: str = "primary", role_b: str = "accent_1",
    role_c: str = "accent_2", role_d: str = "text_muted",
    role_header_text: str = "text_inverse", role_body_text: str = "text_primary",
) -> None:
    """2x2 matrix — classic SWOT/4-quadrant layout.

    quadrants: list of exactly 4 {label, items} dicts in order:
    [top-left, top-right, bottom-left, bottom-right]. items is a list of
    short strings (2-4 each).
    """
    if len(quadrants) < 4:
        # pad if missing
        quadrants = list(quadrants) + [{"label": "", "items": []}] * (4 - len(quadrants))
    quadrants = quadrants[:4]
    cx, cy = anchor_to_emu(anchor)

    grid_w = int(SLIDE_WIDTH_EMU * 0.84)
    grid_h = int(SLIDE_HEIGHT_EMU * 0.72)
    gap = int(SLIDE_WIDTH_EMU * 0.012)
    cell_w = (grid_w - gap) // 2
    cell_h = (grid_h - gap) // 2

    grid_left = cx - grid_w // 2
    grid_top = max(SAFE_MARGIN_EMU,
                   min(cy - grid_h // 2, SLIDE_HEIGHT_EMU - grid_h - SAFE_MARGIN_EMU))

    header_h = int(cell_h * 0.28)
    body_h = cell_h - header_h
    body_size_pt = size_for_tier(ctx, "body")

    role_for_quad = [role_a, role_b, role_c, role_d]

    for idx, quad in enumerate(quadrants):
        col = idx % 2
        row = idx // 2
        cell_left = grid_left + col * (cell_w + gap)
        cell_top = grid_top + row * (cell_h + gap)

        # Card backdrop (light surface) — subtle gradient + shadow for depth
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, cell_left, cell_top, cell_w, cell_h,
        )
        card.line.fill.background()
        sr, sg, sb = hex_to_rgb(hex_for_role(ctx, "surface"))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(sr, sg, sb)
        _polish_node(card, ctx, "surface", light_pct=0.06, dark_pct=-0.04)

        # Header strip — bolder gradient + accent stop on the right edge
        hr_, hg_, hb_ = hex_to_rgb(hex_for_role(ctx, role_for_quad[idx]))
        header = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            cell_left, cell_top, cell_w, header_h,
        )
        header.line.fill.background()
        header.fill.solid()
        header.fill.fore_color.rgb = RGBColor(hr_, hg_, hb_)
        _polish_node(header, ctx, role_for_quad[idx], shadow=False,
                     light_pct=0.20, dark_pct=-0.08)

        _add_text_fitted(
            slide, ctx, text=str(quad.get("label", "")),
            left=cell_left + Emu(80000), top=cell_top + int(header_h * 0.15),
            width=cell_w - Emu(160000), height=int(header_h * 0.75),
            max_tier="h2", min_pt=14, role=role_header_text, bold=True, font_kind="heading",
        )

        # Body items as a small list
        items = quad.get("items", []) or []
        items = [str(x) for x in items[:4]]
        if items:
            inner_pad = Emu(120000)
            body_top = cell_top + header_h + Emu(80000)
            bbox = slide.shapes.add_textbox(
                cell_left + inner_pad, body_top,
                cell_w - 2 * inner_pad, body_h - Emu(160000),
            )
            tf = bbox.text_frame
            tf.word_wrap = True
            br_, bg_, bb_ = hex_to_rgb(hex_for_role(ctx, role_body_text))
            for i, txt in enumerate(items):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.alignment = PP_ALIGN.LEFT
                run = p.add_run()
                run.text = f"• {txt}"
                run.font.name = font_for_kind(ctx, "body")
                run.font.size = Pt(body_size_pt)
                run.font.color.rgb = RGBColor(br_, bg_, bb_)


def pyramid_layers(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center",
    role_top: str = "primary", role_base: str = "accent_1",
    role_text: str = "text_inverse",
) -> None:
    """A stacked pyramid of trapezoid layers labeled by string.

    items: list ordered top → base. 3-5 layers ideal.
    """
    n = max(2, min(5, len(items)))
    items = items[:n]
    x, y = anchor_to_emu(anchor)

    pyramid_h = int(SLIDE_HEIGHT_EMU * 0.55)
    pyramid_w_top = int(SLIDE_WIDTH_EMU * 0.10)
    pyramid_w_bottom = int(SLIDE_WIDTH_EMU * 0.55)
    layer_h = pyramid_h // n
    cx = x

    top_y = max(SAFE_MARGIN_EMU, min(y - pyramid_h // 2,
                                      SLIDE_HEIGHT_EMU - pyramid_h - SAFE_MARGIN_EMU))

    top_r, top_g, top_b = hex_to_rgb(hex_for_role(ctx, role_top))
    base_r, base_g, base_b = hex_to_rgb(hex_for_role(ctx, role_base))

    for i in range(n):
        # Width interpolated
        t_top = i / n
        t_bot = (i + 1) / n
        w_top = int(pyramid_w_top + (pyramid_w_bottom - pyramid_w_top) * t_top)
        w_bot = int(pyramid_w_top + (pyramid_w_bottom - pyramid_w_top) * t_bot)
        left = cx - w_bot // 2
        layer_top = top_y + i * layer_h

        # Trapezoid using a shape — Trapezoid in MSO_SHAPE uses an isosceles
        # trapezoid wider at base. We use bounding box w_bot.
        shape = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, left, layer_top, w_bot, layer_h)
        shape.line.fill.background()
        # Color interpolation top->base
        ratio = i / max(1, n - 1)
        r = int(top_r + (base_r - top_r) * ratio)
        g = int(top_g + (base_g - top_g) * ratio)
        b = int(top_b + (base_b - top_b) * ratio)
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(r, g, b)
        # Drop shadow for layered depth — keep the interpolated color but add lift
        try:
            ox.apply_outer_shadow(shape, blur=50800, distance=19050,
                                  hex_color="#000000", alpha=30000)
        except Exception:
            pass

        # Label
        _add_text_centered(
            slide, ctx, text=items[i],
            left=left, top=layer_top + int(layer_h * 0.20),
            width=w_bot, height=int(layer_h * 0.6),
            tier="body", role=role_text, bold=True, font_kind="heading",
        )
