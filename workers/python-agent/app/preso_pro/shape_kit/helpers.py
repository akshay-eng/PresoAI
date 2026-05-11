"""Shape kit — phase 2 expansion.

Each helper is a registered function the Composer can call. Signatures here
are what the Executor will dispatch to. Args come from the slide_spec JSON
and have already been validated.

Phase 2 set additions:
  Backgrounds: radial_gradient_bg, mesh_gradient_bg, dot_field_bg, diagonal_band_bg
  Decoratives: oversized_letter, half_circle, gradient_strip
  Cards: solid_card, glass_card, outlined_card
  Typography: body_text, gradient_title, bullet_list
  Composites: stats_row, quote_callout, two_column_text
"""

from __future__ import annotations

from typing import Any

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.slide import Slide
from pptx.util import Emu, Pt

from app.preso_pro.planning.slide_spec import DeckContext
from app.preso_pro.shape_kit._util import (
    font_for_kind,
    hex_for_role,
    hex_to_rgb,
    size_for_tier,
)
from app.preso_pro.shape_kit.anchors import (
    SLIDE_HEIGHT_EMU,
    SLIDE_WIDTH_EMU,
    anchor_to_emu,
)
from app.preso_pro.shape_kit import ooxml_helpers as ox

EMU_PER_PT = 12700
SAFE_MARGIN_EMU = int(SLIDE_WIDTH_EMU * 0.05)


# Per-character display width estimate (in EMUs) per 1pt of font size for the
# heading font. Inter Bold at 96pt renders ~10K EMU per char on screen; we use
# 11000 for a small safety margin to avoid horizontal overflow without being
# so conservative that titles end up tiny.
_AVG_CHAR_WIDTH_EMU_PER_PT = 11000

# Per-tier minimum font sizes — a display title should never shrink below
# 48pt no matter what; otherwise it stops being a "display" title.
_TIER_MIN_PT = {
    "display": 48,
    "h1": 32,
    "h2": 24,
    "body": 14,
    "caption": 10,
}


def _fit_font_size(
    text: str,
    max_pt: int,
    max_width_emu: int,
    max_lines: int = 2,
    tier: str = "h1",
) -> int:
    """Largest font size <= max_pt that lets `text` fit in max_width_emu over
    up to `max_lines` lines, but never below the tier's floor."""
    if not text:
        return max_pt
    char_count = max(1, len(text))
    chars_per_line = max(1, char_count / max_lines)
    needed_pt = max_width_emu / (chars_per_line * _AVG_CHAR_WIDTH_EMU_PER_PT)
    floor = _TIER_MIN_PT.get(tier, 18)
    return max(floor, min(max_pt, int(needed_pt)))


def _vertical_top_for_anchor(anchor: str, y: int, height: int) -> int:
    """Resolve the top edge of a box given the anchor's vertical region.

    upper-* anchors -> top edge at y (text grows down)
    lower-* anchors -> bottom edge at y (text grows up)
    center-* anchors / fallback -> centered around y
    """
    if anchor.startswith("upper"):
        return max(0, y - int(0.02 * SLIDE_HEIGHT_EMU))
    if anchor.startswith("lower"):
        return max(0, y - height + int(0.02 * SLIDE_HEIGHT_EMU))
    return max(0, y - height // 2)


def _send_to_back(slide: Slide, shape: Any) -> None:
    sp_tree = slide.shapes._spTree
    sp_tree.remove(shape._element)
    sp_tree.insert(2, shape._element)


def _add_solid_rect(
    slide: Slide,
    left: int,
    top: int,
    width: int,
    height: int,
    hex_color: str,
    no_line: bool = True,
) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if no_line:
        shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_color)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    return shape


# ─── Backgrounds ──────────────────────────────────────────────────────────

def solid_bg(slide: Slide, ctx: DeckContext, *, role: str = "background") -> None:
    """Fill entire slide with a single palette role color."""
    shape = _add_solid_rect(
        slide, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU, hex_for_role(ctx, role)
    )
    _send_to_back(slide, shape)


def linear_gradient_bg(
    slide: Slide, ctx: DeckContext,
    *, role_a: str = "background", role_b: str = "primary", angle: int = 135,
) -> None:
    """Two-stop linear gradient covering the slide."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU
    )
    shape.line.fill.background()
    ox.apply_linear_gradient(shape, hex_for_role(ctx, role_a), hex_for_role(ctx, role_b), angle)
    _send_to_back(slide, shape)


def radial_gradient_bg(
    slide: Slide, ctx: DeckContext,
    *, role_inner: str = "primary", role_outer: str = "background",
    focus: str = "center",
) -> None:
    """Radial gradient. focus is one of: 'center', 'upper-left', 'upper-right',
    'lower-left', 'lower-right'."""
    focus_map = {
        "center": (50000, 50000),
        "upper-left": (15000, 15000),
        "upper-right": (85000, 15000),
        "lower-left": (15000, 85000),
        "lower-right": (85000, 85000),
    }
    fx, fy = focus_map.get(focus, (50000, 50000))
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU
    )
    shape.line.fill.background()
    ox.apply_radial_gradient(
        shape, hex_for_role(ctx, role_inner), hex_for_role(ctx, role_outer),
        focus_x=fx, focus_y=fy,
    )
    _send_to_back(slide, shape)


def mesh_gradient_bg(
    slide: Slide, ctx: DeckContext,
    *, roles: list[str] | None = None, blob_count: int = 3,
) -> None:
    """Approximated mesh gradient: solid base + corner-anchored translucent blobs.

    Blobs are smaller than the slide so content remains legible. They sit at
    corners to create a soft glow rather than dominate the slide.
    """
    roles = roles or ["background", "primary", "accent_1", "accent_2"]
    if not roles:
        roles = ["background"]

    base = _add_solid_rect(
        slide, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU, hex_for_role(ctx, roles[0])
    )
    _send_to_back(slide, base)

    # Corner-positioned, smaller blobs. Sized so they don't cover content.
    # Each: (center_x_frac, center_y_frac, diameter_frac)
    blob_positions = [
        (-0.05, -0.05, 0.35),  # top-left, partly off-slide
        (1.05, -0.05, 0.35),   # top-right, partly off-slide
        (1.10, 1.05, 0.40),    # bottom-right, partly off-slide
        (-0.10, 1.10, 0.35),   # bottom-left, partly off-slide
    ][:max(1, min(4, blob_count))]

    for i, (fx, fy, fsize) in enumerate(blob_positions):
        role = roles[(i + 1) % len(roles)] if len(roles) > 1 else roles[0]
        size = int(SLIDE_WIDTH_EMU * fsize)
        cx = int(fx * SLIDE_WIDTH_EMU)
        cy = int(fy * SLIDE_HEIGHT_EMU)
        left = cx - size // 2
        top = cy - size // 2
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        oval.line.fill.background()
        r, g, b = hex_to_rgb(hex_for_role(ctx, role))
        oval.fill.solid()
        oval.fill.fore_color.rgb = RGBColor(r, g, b)
        ox.apply_alpha(oval, 0.45)
        sp_tree = slide.shapes._spTree
        sp_tree.remove(oval._element)
        sp_tree.insert(3 + i, oval._element)


def dot_field_bg(
    slide: Slide, ctx: DeckContext,
    *, role_bg: str = "background", role_dot: str = "primary",
    density: int = 8,
) -> None:
    """Solid background with a regular grid of small dots in the accent role."""
    base = _add_solid_rect(
        slide, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU, hex_for_role(ctx, role_bg)
    )
    _send_to_back(slide, base)
    cols = max(4, min(16, density))
    rows = max(2, cols // 2)
    cell_w = SLIDE_WIDTH_EMU // cols
    cell_h = SLIDE_HEIGHT_EMU // rows
    dot_size = int(min(cell_w, cell_h) * 0.18)
    r, g, b = hex_to_rgb(hex_for_role(ctx, role_dot))
    for ci in range(cols):
        for ri in range(rows):
            cx = (ci + 0.5) * cell_w
            cy = (ri + 0.5) * cell_h
            left = int(cx - dot_size / 2)
            top = int(cy - dot_size / 2)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, dot_size, dot_size)
            dot.line.fill.background()
            dot.fill.solid()
            dot.fill.fore_color.rgb = RGBColor(r, g, b)
            ox.apply_alpha(dot, 0.18)


def diagonal_band_bg(
    slide: Slide, ctx: DeckContext,
    *, role_bg: str = "background", role_band: str = "primary",
) -> None:
    """Solid background plus a diagonal band stripe (parallelogram across slide)."""
    base = _add_solid_rect(
        slide, 0, 0, SLIDE_WIDTH_EMU, SLIDE_HEIGHT_EMU, hex_for_role(ctx, role_bg)
    )
    _send_to_back(slide, base)
    band_h = int(SLIDE_HEIGHT_EMU * 0.35)
    band = slide.shapes.add_shape(
        MSO_SHAPE.PARALLELOGRAM,
        -int(SLIDE_WIDTH_EMU * 0.10),
        int(SLIDE_HEIGHT_EMU * 0.30),
        int(SLIDE_WIDTH_EMU * 1.20),
        band_h,
    )
    band.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role_band))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(r, g, b)
    ox.apply_alpha(band, 0.20)


# ─── Decorative shapes ─────────────────────────────────────────────────────

def accent_blob(
    slide: Slide, ctx: DeckContext,
    *, position: str = "upper-right", size: int = 200,
    role: str = "accent_1", opacity: float = 0.5,
) -> None:
    """Decorative oval accent with optional transparency."""
    x, y = anchor_to_emu(position)
    side = int(size * EMU_PER_PT)
    left = x - side // 2
    top = y - side // 2

    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, side, side)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    ox.apply_alpha(shape, opacity)


def half_circle(
    slide: Slide, ctx: DeckContext,
    *, position: str = "lower-right", size: int = 250,
    role: str = "accent_2", opacity: float = 0.7,
) -> None:
    """A half-circle decorative — uses the PIE shape."""
    x, y = anchor_to_emu(position)
    side = int(size * EMU_PER_PT)
    left = x - side // 2
    top = y - side // 2
    shape = slide.shapes.add_shape(MSO_SHAPE.PIE, left, top, side, side)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    ox.apply_alpha(shape, opacity)


def gradient_strip(
    slide: Slide, ctx: DeckContext,
    *, role_a: str = "primary", role_b: str = "accent_1",
    edge: str = "left", thickness: int = 60,
) -> None:
    """A thin gradient stripe along one edge of the slide. edge: left|right|top|bottom."""
    th_emu = int(thickness * EMU_PER_PT)
    if edge == "left":
        left, top, w, h = 0, 0, th_emu, SLIDE_HEIGHT_EMU
    elif edge == "right":
        left, top, w, h = SLIDE_WIDTH_EMU - th_emu, 0, th_emu, SLIDE_HEIGHT_EMU
    elif edge == "top":
        left, top, w, h = 0, 0, SLIDE_WIDTH_EMU, th_emu
    else:  # bottom
        left, top, w, h = 0, SLIDE_HEIGHT_EMU - th_emu, SLIDE_WIDTH_EMU, th_emu
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.line.fill.background()
    angle = 90 if edge in ("top", "bottom") else 0
    ox.apply_linear_gradient(shape, hex_for_role(ctx, role_a), hex_for_role(ctx, role_b), angle)


def oversized_letter(
    slide: Slide, ctx: DeckContext,
    *, text: str = "01", anchor: str = "lower-left",
    role: str = "primary", opacity: float = 0.10,
    size: int = 600,
) -> None:
    """Decorative giant letter/number/short string used as visual interest."""
    x, y = anchor_to_emu(anchor)
    width = int(SLIDE_WIDTH_EMU * 0.55)
    height = int(size * EMU_PER_PT * 1.2)
    left = max(-int(SLIDE_WIDTH_EMU * 0.05), x - width // 2)
    top = max(-int(SLIDE_HEIGHT_EMU * 0.05), y - height // 2)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, "heading")
    run.font.size = Pt(min(size, 600))
    run.font.bold = True
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    run.font.color.rgb = RGBColor(r, g, b)
    # opacity via alpha on the text run color
    rPr = run._r.get_or_add_rPr()
    fill = rPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill")
    if fill is not None:
        srgb = fill.find("{http://schemas.openxmlformats.org/drawingml/2006/main}srgbClr")
        if srgb is not None:
            from lxml import etree
            existing_alpha = srgb.find("{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
            if existing_alpha is not None:
                srgb.remove(existing_alpha)
            alpha = etree.SubElement(srgb, "{http://schemas.openxmlformats.org/drawingml/2006/main}alpha")
            alpha.set("val", str(int(opacity * 100000)))


# ─── Cards & containers ────────────────────────────────────────────────────

def solid_card(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.40, height_frac: float = 0.40,
    role: str = "surface", radius: int = 24,
) -> None:
    """A solid rounded rectangle card. Doesn't carry text by itself —
    typically paired with title_text/body_text on top."""
    x, y = anchor_to_emu(anchor)
    w = int(SLIDE_WIDTH_EMU * width_frac)
    h = int(SLIDE_HEIGHT_EMU * height_frac)
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)


def glass_card(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.40, height_frac: float = 0.40,
    role: str = "surface", border_role: str = "accent_1",
    opacity: float = 0.30,
) -> None:
    """Translucent rounded card with a subtle accent border."""
    x, y = anchor_to_emu(anchor)
    w = int(SLIDE_WIDTH_EMU * width_frac)
    h = int(SLIDE_HEIGHT_EMU * height_frac)
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(r, g, b)
    ox.apply_alpha(shape, opacity)
    # thin accent border
    br, bg, bb = hex_to_rgb(hex_for_role(ctx, border_role))
    shape.line.color.rgb = RGBColor(br, bg, bb)
    shape.line.width = Emu(12700)  # 1pt


def outlined_card(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.40, height_frac: float = 0.40,
    border_role: str = "primary", radius: int = 24,
) -> None:
    """Rounded card with no fill — outline only."""
    x, y = anchor_to_emu(anchor)
    w = int(SLIDE_WIDTH_EMU * width_frac)
    h = int(SLIDE_HEIGHT_EMU * height_frac)
    left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, border_role))
    shape.line.color.rgb = RGBColor(r, g, b)
    shape.line.width = Emu(19050)  # 1.5pt


# ─── Typography ────────────────────────────────────────────────────────────

def title_text(
    slide: Slide, ctx: DeckContext,
    *, text: str, anchor: str = "center-left", tier: str = "h1",
    role: str = "text_primary", weight: int = 700, width_frac: float = 0.62,
) -> None:
    """Render a title text block at the named anchor."""
    x, y = anchor_to_emu(anchor)
    width = int(SLIDE_WIDTH_EMU * width_frac)
    base_size = size_for_tier(ctx, tier)
    font_size = _fit_font_size(text, base_size, width, max_lines=2, tier=tier)
    line_h = int(font_size * EMU_PER_PT * 1.4)
    height = line_h * 2

    left = max(SAFE_MARGIN_EMU, x - int(0.05 * SLIDE_WIDTH_EMU))
    if left + width > SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU:
        width = SLIDE_WIDTH_EMU - left - SAFE_MARGIN_EMU
    top = _vertical_top_for_anchor(anchor, y, height)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, "heading")
    run.font.size = Pt(font_size)
    run.font.bold = weight >= 600
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    run.font.color.rgb = RGBColor(r, g, b)


def body_text(
    slide: Slide, ctx: DeckContext,
    *, text: str, anchor: str = "center-left", tier: str = "body",
    role: str = "text_muted", width_frac: float = 0.55,
) -> None:
    """Render regular paragraph body text at an anchor."""
    x, y = anchor_to_emu(anchor)
    font_size = size_for_tier(ctx, tier)
    width = int(SLIDE_WIDTH_EMU * width_frac)
    line_h = int(font_size * EMU_PER_PT * 1.5)
    height = line_h * 6

    left = max(SAFE_MARGIN_EMU, x - int(0.05 * SLIDE_WIDTH_EMU))
    if left + width > SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU:
        width = SLIDE_WIDTH_EMU - left - SAFE_MARGIN_EMU
    top = _vertical_top_for_anchor(anchor, y, height)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, "body")
    run.font.size = Pt(font_size)
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    run.font.color.rgb = RGBColor(r, g, b)


def gradient_title(
    slide: Slide, ctx: DeckContext,
    *, text: str, anchor: str = "center-left", tier: str = "display",
    role_a: str = "primary", role_b: str = "accent_1",
    weight: int = 800, width_frac: float = 0.78,
) -> None:
    """A title with a gradient text fill — visually striking, marketing-grade."""
    x, y = anchor_to_emu(anchor)
    width = int(SLIDE_WIDTH_EMU * width_frac)
    base_size = size_for_tier(ctx, tier)
    font_size = _fit_font_size(text, base_size, width, max_lines=2, tier=tier)
    line_h = int(font_size * EMU_PER_PT * 1.4)
    height = line_h * 2

    left = max(SAFE_MARGIN_EMU, x - int(0.05 * SLIDE_WIDTH_EMU))
    if left + width > SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU:
        width = SLIDE_WIDTH_EMU - left - SAFE_MARGIN_EMU
    top = _vertical_top_for_anchor(anchor, y, height)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, "heading")
    run.font.size = Pt(font_size)
    run.font.bold = weight >= 600
    ox.apply_text_gradient(run, hex_for_role(ctx, role_a), hex_for_role(ctx, role_b), 90)


def bullet_list(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center-left",
    tier: str = "body", role: str = "text_primary",
    width_frac: float = 0.55,
) -> None:
    """A simple bullet list. items is a list of short strings."""
    x, y = anchor_to_emu(anchor)
    font_size = size_for_tier(ctx, tier)
    width = int(SLIDE_WIDTH_EMU * width_frac)
    line_h = int(font_size * EMU_PER_PT * 1.6)
    height = line_h * (len(items) + 1)

    left = max(SAFE_MARGIN_EMU, x - int(0.05 * SLIDE_WIDTH_EMU))
    if left + width > SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU:
        width = SLIDE_WIDTH_EMU - left - SAFE_MARGIN_EMU
    top = _vertical_top_for_anchor(anchor, y, height)

    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))

    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"·  {item}"
        run.font.name = font_for_kind(ctx, "body")
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(r, g, b)
        p.space_after = Emu(60000)


# ─── Composite primitives ─────────────────────────────────────────────────

def data_callout(
    slide: Slide, ctx: DeckContext,
    *, position: str = "lower-right", value: str, label: str,
    role: str = "accent_1",
) -> None:
    """Big-number-with-caption pattern."""
    x, y = anchor_to_emu(position)
    value_size = size_for_tier(ctx, "display")
    label_size = size_for_tier(ctx, "body")

    width = int(SLIDE_WIDTH_EMU * 0.30)
    value_h = int(value_size * EMU_PER_PT * 1.2)
    label_h = int(label_size * EMU_PER_PT * 1.4)
    total_h = value_h + label_h + Emu(60000)

    left = max(SAFE_MARGIN_EMU, min(x - width // 2, SLIDE_WIDTH_EMU - width - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - total_h // 2, SLIDE_HEIGHT_EMU - total_h - SAFE_MARGIN_EMU))

    value_box = slide.shapes.add_textbox(left, top, width, value_h)
    vtf = value_box.text_frame
    vtf.word_wrap = True
    vp = vtf.paragraphs[0]
    vrun = vp.add_run()
    vrun.text = value
    vrun.font.name = font_for_kind(ctx, "heading")
    vrun.font.size = Pt(value_size)
    vrun.font.bold = True
    vr, vg, vb = hex_to_rgb(hex_for_role(ctx, role))
    vrun.font.color.rgb = RGBColor(vr, vg, vb)

    label_top = top + value_h + 60000
    label_box = slide.shapes.add_textbox(left, label_top, width, label_h)
    ltf = label_box.text_frame
    ltf.word_wrap = True
    lp = ltf.paragraphs[0]
    lrun = lp.add_run()
    lrun.text = label
    lrun.font.name = font_for_kind(ctx, "body")
    lrun.font.size = Pt(label_size)
    lr, lg, lb = hex_to_rgb(hex_for_role(ctx, "text_muted"))
    lrun.font.color.rgb = RGBColor(lr, lg, lb)


def stats_row(
    slide: Slide, ctx: DeckContext,
    *, items: list[dict], anchor: str = "center", role: str = "primary",
) -> None:
    """A horizontal row of N stat tiles. items: [{value, label}, ...] (2-5 items).

    Renders a centered row across the slide width with even spacing."""
    n = max(1, min(5, len(items)))
    items = items[:n]
    x, y = anchor_to_emu(anchor)

    # Pick stat-value type tier so 5 cells don't overflow (e.g. "$1.2M" needs room).
    if n >= 4:
        value_tier = "h2"
    else:
        value_tier = "h1"
    value_size = size_for_tier(ctx, value_tier)
    label_size = size_for_tier(ctx, "caption")
    total_w = int(SLIDE_WIDTH_EMU * 0.84)
    cell_w = total_w // n
    value_h = int(value_size * EMU_PER_PT * 1.2)
    label_h = int(label_size * EMU_PER_PT * 1.6)
    total_h = value_h + label_h + Emu(60000)

    start_left = (SLIDE_WIDTH_EMU - total_w) // 2
    base_top = _vertical_top_for_anchor(anchor, y, total_h)
    top = max(SAFE_MARGIN_EMU, min(base_top, SLIDE_HEIGHT_EMU - total_h - SAFE_MARGIN_EMU))

    for i, item in enumerate(items):
        cell_left = start_left + i * cell_w
        # value
        vbox = slide.shapes.add_textbox(cell_left, top, cell_w, value_h)
        vtf = vbox.text_frame
        vtf.word_wrap = False
        vp = vtf.paragraphs[0]
        vp.alignment = PP_ALIGN.CENTER
        vrun = vp.add_run()
        vrun.text = str(item.get("value", ""))
        vrun.font.name = font_for_kind(ctx, "heading")
        vrun.font.size = Pt(value_size)
        vrun.font.bold = True
        vr, vg, vb = hex_to_rgb(hex_for_role(ctx, role))
        vrun.font.color.rgb = RGBColor(vr, vg, vb)

        # label
        lbox = slide.shapes.add_textbox(cell_left, top + value_h + 60000, cell_w, label_h)
        ltf = lbox.text_frame
        ltf.word_wrap = True
        lp = ltf.paragraphs[0]
        lp.alignment = PP_ALIGN.CENTER
        lrun = lp.add_run()
        lrun.text = str(item.get("label", ""))
        lrun.font.name = font_for_kind(ctx, "body")
        lrun.font.size = Pt(label_size)
        lr, lg, lb = hex_to_rgb(hex_for_role(ctx, "text_muted"))
        lrun.font.color.rgb = RGBColor(lr, lg, lb)


def quote_callout(
    slide: Slide, ctx: DeckContext,
    *, text: str, attribution: str = "",
    anchor: str = "center", role_quote: str = "text_primary",
    role_accent: str = "accent_1",
) -> None:
    """Pull-quote pattern: large italicized text plus attribution line below."""
    x, y = anchor_to_emu(anchor)
    quote_size = size_for_tier(ctx, "h2")
    attr_size = size_for_tier(ctx, "body")
    width = int(SLIDE_WIDTH_EMU * 0.70)
    quote_lines = max(2, len(text) // 60 + 1)
    quote_h = int(quote_size * EMU_PER_PT * 1.4 * quote_lines)
    attr_h = int(attr_size * EMU_PER_PT * 1.4)
    total_h = quote_h + attr_h + Emu(60000)

    left = (SLIDE_WIDTH_EMU - width) // 2
    top = max(SAFE_MARGIN_EMU, min(y - total_h // 2, SLIDE_HEIGHT_EMU - total_h - SAFE_MARGIN_EMU))

    qbox = slide.shapes.add_textbox(left, top, width, quote_h)
    qtf = qbox.text_frame
    qtf.word_wrap = True
    qp = qtf.paragraphs[0]
    qrun = qp.add_run()
    qrun.text = f"“{text}”"
    qrun.font.name = font_for_kind(ctx, "heading")
    qrun.font.size = Pt(quote_size)
    qrun.font.italic = True
    qr, qg, qb = hex_to_rgb(hex_for_role(ctx, role_quote))
    qrun.font.color.rgb = RGBColor(qr, qg, qb)

    if attribution:
        abox = slide.shapes.add_textbox(left, top + quote_h + 60000, width, attr_h)
        atf = abox.text_frame
        ap = atf.paragraphs[0]
        arun = ap.add_run()
        arun.text = f"— {attribution}"
        arun.font.name = font_for_kind(ctx, "body")
        arun.font.size = Pt(attr_size)
        ar, ag, ab = hex_to_rgb(hex_for_role(ctx, role_accent))
        arun.font.color.rgb = RGBColor(ar, ag, ab)


def card_stack(
    slide: Slide, ctx: DeckContext,
    *, anchor: str = "center", width_frac: float = 0.55, height_frac: float = 0.50,
    role: str = "surface", accent_role: str = "accent_1", count: int = 2,
) -> None:
    """Visually rich: 2-3 offset rounded cards layered behind each other.
    Foreground card is opaque (`role`); background offset card is translucent
    accent. Used as a base for hero/quote slides for depth."""
    x, y = anchor_to_emu(anchor)
    w = int(SLIDE_WIDTH_EMU * width_frac)
    h = int(SLIDE_HEIGHT_EMU * height_frac)
    base_left = max(SAFE_MARGIN_EMU, min(x - w // 2, SLIDE_WIDTH_EMU - w - SAFE_MARGIN_EMU))
    base_top = max(SAFE_MARGIN_EMU, min(y - h // 2, SLIDE_HEIGHT_EMU - h - SAFE_MARGIN_EMU))

    n = max(1, min(3, count))
    offsets = [
        (int(SLIDE_WIDTH_EMU * 0.025), int(SLIDE_HEIGHT_EMU * 0.030), accent_role, 0.55),
        (int(SLIDE_WIDTH_EMU * -0.020), int(SLIDE_HEIGHT_EMU * -0.025), "primary", 0.40),
    ]
    # background cards first, foreground last
    for i in range(n - 1):
        ox_off, oy_off, role_off, alpha_off = offsets[i % len(offsets)]
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            base_left + ox_off, base_top + oy_off, w, h,
        )
        card.line.fill.background()
        r, g, b = hex_to_rgb(hex_for_role(ctx, role_off))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(r, g, b)
        ox.apply_alpha(card, alpha_off)
    # foreground (solid surface)
    fg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, base_left, base_top, w, h)
    fg.line.fill.background()
    r, g, b = hex_to_rgb(hex_for_role(ctx, role))
    fg.fill.solid()
    fg.fill.fore_color.rgb = RGBColor(r, g, b)
    ox.apply_outer_shadow(fg, blur=63500, distance=25400, alpha=40000)


def glass_text_panel(
    slide: Slide, ctx: DeckContext,
    *, text: str, anchor: str = "center", tier: str = "h1",
    role_text: str = "text_primary", role_panel: str = "surface",
    role_border: str = "accent_1", width_frac: float = 0.55, padding_frac: float = 0.04,
) -> None:
    """Translucent panel with text inside — modern marketing 'card with text' pattern."""
    x, y = anchor_to_emu(anchor)
    width = int(SLIDE_WIDTH_EMU * width_frac)
    base_size = size_for_tier(ctx, tier)
    font_size = _fit_font_size(text, base_size, int(width * (1 - 2 * padding_frac)), max_lines=3)
    line_h = int(font_size * EMU_PER_PT * 1.5)
    inner_h = line_h * 3
    pad = int(SLIDE_WIDTH_EMU * padding_frac)
    height = inner_h + pad * 2

    left = max(SAFE_MARGIN_EMU, min(x - width // 2, SLIDE_WIDTH_EMU - width - SAFE_MARGIN_EMU))
    top = max(SAFE_MARGIN_EMU, min(y - height // 2, SLIDE_HEIGHT_EMU - height - SAFE_MARGIN_EMU))

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pr, pg, pb = hex_to_rgb(hex_for_role(ctx, role_panel))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(pr, pg, pb)
    ox.apply_alpha(panel, 0.30)
    br, bg, bb = hex_to_rgb(hex_for_role(ctx, role_border))
    panel.line.color.rgb = RGBColor(br, bg, bb)
    panel.line.width = Emu(12700)

    box = slide.shapes.add_textbox(left + pad, top + pad, width - 2 * pad, inner_h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = font_for_kind(ctx, "heading")
    run.font.size = Pt(font_size)
    run.font.bold = True
    tr, tg, tb = hex_to_rgb(hex_for_role(ctx, role_text))
    run.font.color.rgb = RGBColor(tr, tg, tb)


def corner_accent_cluster(
    slide: Slide, ctx: DeckContext,
    *, corner: str = "lower-right", role: str = "accent_1",
    secondary_role: str = "accent_2", scale: float = 1.0,
) -> None:
    """A grouping of small geometric shapes anchored at a corner — adds energy
    to a slide without dominating it. Combines a circle, a small rounded
    rectangle, and a triangle."""
    base = int(SLIDE_WIDTH_EMU * 0.06 * scale)
    # corner-coord generators; (origin_x, origin_y, vector_x, vector_y)
    corner_origins = {
        "upper-left": (SAFE_MARGIN_EMU, SAFE_MARGIN_EMU, 1, 1),
        "upper-right": (SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU, SAFE_MARGIN_EMU, -1, 1),
        "lower-left": (SAFE_MARGIN_EMU, SLIDE_HEIGHT_EMU - SAFE_MARGIN_EMU, 1, -1),
        "lower-right": (SLIDE_WIDTH_EMU - SAFE_MARGIN_EMU, SLIDE_HEIGHT_EMU - SAFE_MARGIN_EMU, -1, -1),
    }
    ox_, oy_, vx, vy = corner_origins.get(corner, corner_origins["lower-right"])

    pieces = [
        (MSO_SHAPE.OVAL, base, base, role, 0.85, 0, 0),
        (MSO_SHAPE.ROUNDED_RECTANGLE, int(base * 1.6), int(base * 0.5), secondary_role, 0.55,
         int(base * 1.5) * vx, int(base * 0.4) * vy),
        (MSO_SHAPE.RIGHT_TRIANGLE, int(base * 0.7), int(base * 0.7), role, 0.40,
         int(base * 0.4) * vx, int(base * 1.3) * vy),
    ]
    for shape_kind, w, h, prole, alpha, dx, dy in pieces:
        # anchor: shape's far corner sits at (ox + dx) etc.
        if vx > 0:
            left = ox_ + dx
        else:
            left = ox_ + dx - w
        if vy > 0:
            top = oy_ + dy
        else:
            top = oy_ + dy - h
        s = slide.shapes.add_shape(shape_kind, left, top, w, h)
        s.line.fill.background()
        r, g, b = hex_to_rgb(hex_for_role(ctx, prole))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(r, g, b)
        ox.apply_alpha(s, alpha)


def oversized_number(
    slide: Slide, ctx: DeckContext,
    *, text: str = "01", anchor: str = "lower-right",
    role: str = "accent_1", opacity: float = 0.20, size: int = 380,
) -> None:
    """Like oversized_letter but specifically tuned for stat/number decoration.
    Smaller default size, more saturated default, designed to peek behind content."""
    oversized_letter(
        slide, ctx,
        text=text, anchor=anchor, role=role, opacity=opacity, size=size,
    )


def two_column_text(
    slide: Slide, ctx: DeckContext,
    *, left_title: str = "", left_body: str = "",
    right_title: str = "", right_body: str = "",
    anchor: str = "center",
    role_title: str = "text_primary", role_body: str = "text_muted",
) -> None:
    """Two-column layout: title + body on each side."""
    x, y = anchor_to_emu(anchor)
    col_width = int(SLIDE_WIDTH_EMU * 0.40)
    gap = int(SLIDE_WIDTH_EMU * 0.05)
    total_w = col_width * 2 + gap
    title_size = size_for_tier(ctx, "h2")
    body_size = size_for_tier(ctx, "body")
    title_h = int(title_size * EMU_PER_PT * 1.4)
    body_h = int(body_size * EMU_PER_PT * 1.6 * 5)
    total_h = title_h + body_h + Emu(60000)

    start_left = (SLIDE_WIDTH_EMU - total_w) // 2
    top = max(SAFE_MARGIN_EMU, min(y - total_h // 2, SLIDE_HEIGHT_EMU - total_h - SAFE_MARGIN_EMU))

    columns = [
        (start_left, left_title, left_body),
        (start_left + col_width + gap, right_title, right_body),
    ]
    for col_left, ctitle, cbody in columns:
        if ctitle:
            tbox = slide.shapes.add_textbox(col_left, top, col_width, title_h)
            ttf = tbox.text_frame
            ttf.word_wrap = True
            tp = ttf.paragraphs[0]
            tr = tp.add_run()
            tr.text = ctitle
            tr.font.name = font_for_kind(ctx, "heading")
            tr.font.size = Pt(title_size)
            tr.font.bold = True
            r, g, b = hex_to_rgb(hex_for_role(ctx, role_title))
            tr.font.color.rgb = RGBColor(r, g, b)
        if cbody:
            bbox = slide.shapes.add_textbox(col_left, top + title_h + 60000, col_width, body_h)
            btf = bbox.text_frame
            btf.word_wrap = True
            bp = btf.paragraphs[0]
            br = bp.add_run()
            br.text = cbody
            br.font.name = font_for_kind(ctx, "body")
            br.font.size = Pt(body_size)
            r, g, b = hex_to_rgb(hex_for_role(ctx, role_body))
            br.font.color.rgb = RGBColor(r, g, b)
