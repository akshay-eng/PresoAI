"""Shape kit functions for real native SmartArt — visually rich layout set.

Each function below:
1. Drops a marker rectangle on the slide so the post-processor knows where the
   SmartArt should land.
2. Records a PendingSmartArt request on the presentation's `_pending_smart_art`
   list. The orchestrator drains this list after python-pptx saves the file.

The layouts here are the curated set the user picked from their reference
file — modern, geometric, with real shape backgrounds (not just text + thin
lines). LibreOffice/Collabora renders these well; PowerPoint renders them at
full quality.

Slot counts (max items the v1 template supports per call):
  cycle3                          5  — Continuous Cycle (5 boxes with arc arrows)
  orgChart1                       4  — root + 3 children
  lProcess3                       9  — 3 rows × 3 chevrons
  lProcess2                       9  — 3 columns × 3 cells (title + 2 items each)
  hierarchy5                      9  — 3-column hierarchy
  HexagonTimeline                 6  — 3 hexagons × (title + caption)
  AccentHomeChevronProcess        6  — 3 chevrons × (title + caption)
  hList6                          9  — 3 trapezoid columns × 3 text levels
"""

from __future__ import annotations

from pptx.slide import Slide

from app.preso_pro.planning.slide_spec import DeckContext
from app.preso_pro.shape_kit.anchors import (
    SLIDE_HEIGHT_EMU,
    SLIDE_WIDTH_EMU,
    anchor_to_emu,
)
from app.preso_pro.smart_art.builder import (
    SmartArtBuilder,
    SmartArtRequest,
    new_marker_id,
)
from app.preso_pro.smart_art.package_helpers import MARKER_PREFIX

SAFE_MARGIN_EMU = int(SLIDE_WIDTH_EMU * 0.05)


def _frame_geometry(anchor: str, size_frac: float) -> tuple[int, int, int, int]:
    """Compute (x, y, cx, cy) EMU for the SmartArt frame, centered on anchor."""
    cx = int(SLIDE_WIDTH_EMU * size_frac)
    cy = int(SLIDE_HEIGHT_EMU * size_frac)
    ax, ay = anchor_to_emu(anchor)
    x = max(SAFE_MARGIN_EMU, min(ax - cx // 2, SLIDE_WIDTH_EMU - cx - SAFE_MARGIN_EMU))
    y = max(SAFE_MARGIN_EMU, min(ay - cy // 2, SLIDE_HEIGHT_EMU - cy - SAFE_MARGIN_EMU))
    return x, y, cx, cy


def _drop_marker(slide: Slide, marker_id: str, x: int, y: int, cx: int, cy: int) -> None:
    """Add a hidden rectangle whose name encodes the marker_id. The
    post-processor finds it by name and replaces it with a graphicFrame."""
    from pptx.enum.shapes import MSO_SHAPE
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, cx, cy)
    shape.name = MARKER_PREFIX + marker_id
    shape.line.fill.background()
    shape.fill.background()


def _enqueue(slide: Slide, layout_key: str, items: list[str],
             anchor: str, size_frac: float) -> None:
    """Common path: place marker, build pending request, push to package list."""
    x, y, cx, cy = _frame_geometry(anchor, size_frac)
    marker_id = new_marker_id()

    package = slide.part.package
    pending_list = getattr(package, "_pending_smart_art", None)
    if pending_list is None:
        pending_list = []
        package._pending_smart_art = pending_list  # type: ignore[attr-defined]

    builder = SmartArtBuilder(layout_key)
    pres = package.presentation_part.presentation
    slide_idx = list(pres.slides).index(slide) + 1

    req = SmartArtRequest(
        layout_key=layout_key,
        items=items,
        slide_index=slide_idx,
        x_emu=x, y_emu=y, cx_emu=cx, cy_emu=cy,
    )
    pending_list.append(builder.build(req, marker_id))

    _drop_marker(slide, marker_id, x, y, cx, cy)


# ─── Public LLM-facing functions ──────────────────────────────────────────

def smart_art_cycle(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.78,
) -> None:
    """Native PowerPoint Continuous Cycle SmartArt — 5 boxes around an arc with
    arrows. Visually richer than Basic Cycle. Up to 5 items."""
    _enqueue(slide, "cycle3", items, anchor, size_frac)


def smart_art_org_chart(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.85,
) -> None:
    """Native PowerPoint Organization Chart. Up to 4 items (root + 3 children)."""
    _enqueue(slide, "orgChart1", items, anchor, size_frac)


def smart_art_chevron_rows(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.85,
) -> None:
    """Three rows of three chevron arrows each — for parallel/multi-stream
    processes. 9 items: row1step1, row1step2, row1step3, row2step1, ..."""
    _enqueue(slide, "lProcess3", items, anchor, size_frac)


def smart_art_column_blocks(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.78,
) -> None:
    """Three columns each with a title and stacked cells. Up to 9 items
    (3 columns × 3 levels: column-title, item-1, item-2)."""
    _enqueue(slide, "lProcess2", items, anchor, size_frac)


def smart_art_horizontal_hierarchy(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.85,
) -> None:
    """Horizontal hierarchy with 3 column backgrounds and connector lines.
    Up to 9 items (root + 2 mid-tier + 4 leaves laid out in tree form)."""
    _enqueue(slide, "hierarchy5", items, anchor, size_frac)


def smart_art_hexagon_timeline(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.85,
) -> None:
    """Modern hexagon timeline — 3 hexagons connected by lines, each with a
    title and caption. Up to 6 items (3 milestones × title + caption)."""
    _enqueue(slide, "HexagonTimeline", items, anchor, size_frac)


def smart_art_chevron_process(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.85,
) -> None:
    """Modern accent chevron process — 3 chevrons each with title + caption.
    Up to 6 items."""
    _enqueue(slide, "AccentHomeChevronProcess", items, anchor, size_frac)


def smart_art_trapezoid_blocks(
    slide: Slide, ctx: DeckContext,
    *, items: list[str], anchor: str = "center", size_frac: float = 0.62,
) -> None:
    """Three large trapezoid blocks side-by-side. Bold marketing-style.
    Up to 9 items (3 trapezoids × 3 text levels each)."""
    _enqueue(slide, "hList6", items, anchor, size_frac)
