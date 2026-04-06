from __future__ import annotations

import os
from pathlib import Path
from lxml import etree

import structlog
from pptx import Presentation
from pptx.util import Emu

from app.models import ThemeConfig, ThemeColors, LayoutInfo, PlaceholderInfo
from app.services.s3 import S3Service

logger = structlog.get_logger()

NSMAP = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
}

COLOR_ELEMENT_NAMES = [
    "dk1", "lt1", "dk2", "lt2",
    "accent1", "accent2", "accent3", "accent4", "accent5", "accent6",
    "hlink", "folHlink",
]


class ThemeExtractor:
    def __init__(self) -> None:
        self.s3 = S3Service()

    def extract(self, s3_key: str) -> ThemeConfig:
        tmp_path = self.s3.download_to_temp(s3_key, suffix=".pptx")
        try:
            return self._extract_from_file(tmp_path)
        finally:
            os.unlink(tmp_path)

    def _extract_from_file(self, path: Path) -> ThemeConfig:
        prs = Presentation(str(path))
        colors = self._extract_colors(prs)
        heading_font, body_font = self._extract_fonts(prs)
        layouts = self._extract_layouts(prs)
        master_bg = self._extract_master_background(prs)

        return ThemeConfig(
            colors=colors,
            heading_font=heading_font,
            body_font=body_font,
            layouts=layouts,
            master_background=master_bg,
        )

    def _get_theme_xml(self, prs: Presentation) -> etree._Element | None:
        """Get the theme XML element from the presentation, handling different python-pptx versions."""
        # Method 1: Try slide master element directly
        try:
            slide_master = prs.slide_masters[0]
            clr_scheme = slide_master.element.findall(".//a:clrScheme", NSMAP)
            if clr_scheme:
                return slide_master.element
        except Exception:
            pass

        # Method 2: Try iter_parts (newer python-pptx)
        try:
            for part in prs.part.package.iter_parts():
                if hasattr(part, 'content_type') and 'theme' in str(part.content_type):
                    return etree.fromstring(part.blob)
        except AttributeError:
            pass

        # Method 3: Try .parts (older python-pptx)
        try:
            for part in prs.part.package.parts:
                if hasattr(part, 'content_type') and 'theme' in str(part.content_type):
                    return etree.fromstring(part.blob)
        except (AttributeError, TypeError):
            pass

        # Method 4: Walk the rels to find theme part
        try:
            master_part = prs.slide_masters[0].part
            for rel in master_part.rels.values():
                if 'theme' in rel.reltype:
                    return etree.fromstring(rel.target_part.blob)
        except Exception:
            pass

        return None

    def _extract_colors(self, prs: Presentation) -> ThemeColors:
        color_map: dict[str, str] = {}
        theme_el = self._get_theme_xml(prs)

        if theme_el is not None:
            clr_scheme_elements = theme_el.findall(".//a:clrScheme", NSMAP)
            if clr_scheme_elements:
                clr_scheme = clr_scheme_elements[0]
                for name in COLOR_ELEMENT_NAMES:
                    el = clr_scheme.find(f"a:{name}", NSMAP)
                    if el is not None:
                        srgb = el.find("a:srgbClr", NSMAP)
                        if srgb is not None:
                            color_map[name] = f"#{srgb.get('val', '000000')}"
                        else:
                            sys_clr = el.find("a:sysClr", NSMAP)
                            if sys_clr is not None:
                                last_clr = sys_clr.get("lastClr", "000000")
                                color_map[name] = f"#{last_clr}"

        return ThemeColors(**color_map) if color_map else ThemeColors()

    def _extract_fonts(self, prs: Presentation) -> tuple[str, str]:
        heading_font = "Calibri"
        body_font = "Calibri"

        theme_el = self._get_theme_xml(prs)
        if theme_el is not None:
            major_font = theme_el.find(".//a:majorFont/a:latin", NSMAP)
            minor_font = theme_el.find(".//a:minorFont/a:latin", NSMAP)
            if major_font is not None:
                heading_font = major_font.get("typeface", "Calibri")
            if minor_font is not None:
                body_font = minor_font.get("typeface", "Calibri")

        return heading_font, body_font

    def _extract_layouts(self, prs: Presentation) -> list[LayoutInfo]:
        layouts: list[LayoutInfo] = []
        for layout in prs.slide_layouts:
            placeholders: list[PlaceholderInfo] = []
            for ph in layout.placeholders:
                try:
                    placeholders.append(
                        PlaceholderInfo(
                            idx=ph.placeholder_format.idx,
                            type=str(ph.placeholder_format.type),
                            x=round(Emu(ph.left).inches, 2) if ph.left else 0,
                            y=round(Emu(ph.top).inches, 2) if ph.top else 0,
                            w=round(Emu(ph.width).inches, 2) if ph.width else 0,
                            h=round(Emu(ph.height).inches, 2) if ph.height else 0,
                        )
                    )
                except Exception:
                    continue
            layouts.append(
                LayoutInfo(
                    name=layout.name,
                    type=self._classify_layout(layout.name, placeholders),
                    placeholders=placeholders,
                )
            )
        return layouts

    @staticmethod
    def _classify_layout(name: str, placeholders: list[PlaceholderInfo]) -> str:
        name_lower = name.lower()
        if "title" in name_lower and ("only" in name_lower or len(placeholders) <= 2):
            return "title"
        if "two" in name_lower or "column" in name_lower:
            return "two_column"
        if "chart" in name_lower:
            return "chart"
        if "image" in name_lower or "picture" in name_lower:
            return "image_focus"
        if "blank" in name_lower:
            return "blank"
        if "section" in name_lower:
            return "section_break"
        return "content"

    def _extract_master_background(self, prs: Presentation) -> str | None:
        try:
            bg = prs.slide_masters[0].background
            if bg and bg.fill and bg.fill.type is not None:
                fore_color = bg.fill.fore_color
                if fore_color and fore_color.rgb:
                    return f"#{fore_color.rgb}"
        except Exception:
            pass
        return None


class ReferenceExtractor:
    def __init__(self) -> None:
        self.s3 = S3Service()

    def extract(self, s3_key: str, file_type: str) -> tuple[str, list[dict]]:
        tmp_path = self.s3.download_to_temp(s3_key)
        try:
            match file_type.lower():
                case "pptx" | ".pptx" | "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                    return self._extract_pptx(tmp_path)
                case "pdf" | ".pdf" | "application/pdf":
                    return self._extract_pdf(tmp_path)
                case "docx" | ".docx" | "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                    return self._extract_docx(tmp_path)
                case _:
                    raise ValueError(f"Unsupported file type: {file_type}")
        finally:
            os.unlink(tmp_path)

    def _extract_pptx(self, path: Path) -> tuple[str, list[dict]]:
        prs = Presentation(str(path))
        full_text_parts: list[str] = []
        structure: list[dict] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_title = ""
            slide_text_parts: list[str] = []

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_text_parts.append(text)
                        try:
                            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                                slide_title = text
                        except Exception:
                            pass

            slide_text = "\n".join(slide_text_parts)
            if slide_text:
                full_text_parts.append(f"--- Slide {i}: {slide_title} ---\n{slide_text}")
                structure.append({
                    "slide_number": i,
                    "title": slide_title,
                    "content": slide_text,
                })

        return "\n\n".join(full_text_parts), structure

    def _extract_pdf(self, path: Path) -> tuple[str, list[dict]]:
        import pdfplumber

        full_text_parts: list[str] = []
        structure: list[dict] = []

        with pdfplumber.open(str(path)) as pdf:
            for i, page in enumerate(pdf.pages, 1):
                text = page.extract_text() or ""
                if text.strip():
                    full_text_parts.append(f"--- Page {i} ---\n{text}")
                    structure.append({
                        "page_number": i,
                        "content": text.strip(),
                    })

        return "\n\n".join(full_text_parts), structure

    def _extract_docx(self, path: Path) -> tuple[str, list[dict]]:
        from docx import Document

        doc = Document(str(path))
        full_text_parts: list[str] = []
        structure: list[dict] = []
        current_section: dict = {"heading": "", "content_parts": []}

        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue

            if para.style and para.style.name and para.style.name.startswith("Heading"):
                if current_section["content_parts"]:
                    content = "\n".join(current_section["content_parts"])
                    full_text_parts.append(
                        f"--- {current_section['heading'] or 'Section'} ---\n{content}"
                    )
                    structure.append({
                        "heading": current_section["heading"],
                        "content": content,
                    })
                current_section = {"heading": text, "content_parts": []}
            else:
                current_section["content_parts"].append(text)

        if current_section["content_parts"]:
            content = "\n".join(current_section["content_parts"])
            full_text_parts.append(
                f"--- {current_section['heading'] or 'Section'} ---\n{content}"
            )
            structure.append({
                "heading": current_section["heading"],
                "content": content,
            })

        return "\n\n".join(full_text_parts), structure
