"""Slide-level indexing: parse PPTX, render slides, OCR, embed (text + image), store with pgvector."""
from __future__ import annotations

import io
import os
import subprocess
import tempfile
from collections import Counter
from pathlib import Path
from typing import List, Optional

import numpy as np
import psycopg
import pytesseract
import structlog
import torch
from PIL import Image
from pgvector.psycopg import register_vector
from pptx import Presentation

from app.config import settings
from app.services.s3 import S3Service

logger = structlog.get_logger()

# Singletons — loaded lazily on first use, kept in memory.
_text_model = None
_clip_model = None
_clip_preprocess = None
_clip_tokenizer = None
_torch_device: Optional[str] = None


def _device() -> str:
    global _torch_device
    if _torch_device is not None:
        return _torch_device
    if torch.backends.mps.is_available():
        _torch_device = "mps"
    elif torch.cuda.is_available():
        _torch_device = "cuda"
    else:
        _torch_device = "cpu"
    logger.info("torch_device_selected", device=_torch_device)
    return _torch_device


def get_text_model():
    global _text_model
    if _text_model is None:
        from sentence_transformers import SentenceTransformer
        logger.info("loading_text_model", name="BAAI/bge-base-en-v1.5")
        _text_model = SentenceTransformer("BAAI/bge-base-en-v1.5", device=_device())
    return _text_model


def get_clip():
    global _clip_model, _clip_preprocess, _clip_tokenizer
    if _clip_model is None:
        import open_clip
        logger.info("loading_clip_model", name="ViT-B-32")
        model, _, preprocess = open_clip.create_model_and_transforms(
            "ViT-B-32", pretrained="laion2b_s34b_b79k"
        )
        model = model.to(_device()).eval()
        _clip_model = model
        _clip_preprocess = preprocess
        _clip_tokenizer = open_clip.get_tokenizer("ViT-B-32")
    return _clip_model, _clip_preprocess, _clip_tokenizer


def embed_text(texts: List[str]) -> np.ndarray:
    """Returns (N, 768) float32 array, L2-normalized for cosine similarity."""
    if not texts:
        return np.zeros((0, 768), dtype=np.float32)
    model = get_text_model()
    # BGE benefits from a query-style prefix only on the search side; we leave it bare for documents.
    embeds = model.encode(texts, normalize_embeddings=True, convert_to_numpy=True, show_progress_bar=False)
    return embeds.astype(np.float32)


def embed_image(images: List[Image.Image]) -> np.ndarray:
    """Returns (N, 512) float32 array, L2-normalized."""
    if not images:
        return np.zeros((0, 512), dtype=np.float32)
    model, preprocess, _ = get_clip()
    device = _device()
    tensors = torch.stack([preprocess(img) for img in images]).to(device)
    with torch.no_grad():
        feats = model.encode_image(tensors)
        feats = feats / feats.norm(dim=-1, keepdim=True)
    return feats.cpu().numpy().astype(np.float32)


def embed_clip_text(text: str) -> np.ndarray:
    """Returns (512,) float32 array — for visual queries against image_embedding."""
    model, _, tokenizer = get_clip()
    device = _device()
    tokens = tokenizer([text]).to(device)
    with torch.no_grad():
        feats = model.encode_text(tokens)
        feats = feats / feats.norm(dim=-1, keepdim=True)
    return feats[0].cpu().numpy().astype(np.float32)


# ── PPTX → slide text ──────────────────────────────────────────────────────

def extract_slide_texts(pptx_path: Path) -> List[str]:
    prs = Presentation(str(pptx_path))
    out: List[str] = []
    for slide in prs.slides:
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.text and run.text.strip():
                            parts.append(run.text.strip())
        # Notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            note = slide.notes_slide.notes_text_frame.text
            if note and note.strip():
                parts.append(note.strip())
        out.append(" ".join(parts))
    return out


# ── PPTX → PNG renders via LibreOffice ────────────────────────────────────

def render_pptx_to_pngs(pptx_path: Path, out_dir: Path) -> List[Path]:
    """Convert PPTX to one PNG per slide using LibreOffice headless."""
    out_dir.mkdir(parents=True, exist_ok=True)
    cmd = [
        settings.libreoffice_path,
        "--headless",
        "--convert-to", "png",
        "--outdir", str(out_dir),
        str(pptx_path),
    ]
    env = os.environ.copy()
    env["HOME"] = str(out_dir)  # LibreOffice writes profile here; isolates state
    logger.info("running_libreoffice", cmd=" ".join(cmd))
    proc = subprocess.run(cmd, capture_output=True, text=True, timeout=180, env=env)
    if proc.returncode != 0:
        logger.error("libreoffice_failed", stderr=proc.stderr[:500])
        raise RuntimeError(f"LibreOffice conversion failed: {proc.stderr[:200]}")
    pngs = sorted(out_dir.glob("*.png"), key=_slide_sort_key)
    return pngs


def _slide_sort_key(p: Path) -> int:
    # LibreOffice typically outputs {basename}.png for single slide and {basename}-{n}.png for multi-slide.
    # Some versions output {basename}_001.png etc. Pull the last integer in stem.
    import re
    nums = re.findall(r"\d+", p.stem)
    return int(nums[-1]) if nums else 0


# ── OCR ───────────────────────────────────────────────────────────────────

def ocr_image(img: Image.Image) -> str:
    try:
        text = pytesseract.image_to_string(img)
        return " ".join(text.split())
    except Exception as exc:
        logger.warning("ocr_failed", error=str(exc))
        return ""


# ── Dominant colors ───────────────────────────────────────────────────────

def dominant_colors(img: Image.Image, k: int = 5) -> list[dict]:
    """Quick 6x6x6 color cube quantization. Returns top-k buckets with hex+weight."""
    small = img.convert("RGB").resize((96, 96))
    arr = np.array(small).reshape(-1, 3)
    # Bin to 6 levels per channel → 216 buckets
    binned = (arr // 43).clip(0, 5)
    keys = binned[:, 0] * 36 + binned[:, 1] * 6 + binned[:, 2]
    counts = Counter(keys.tolist())
    total = sum(counts.values())
    top = counts.most_common(k)
    out = []
    for key, count in top:
        b = key % 6
        g = (key // 6) % 6
        r = key // 36
        # Bucket center in 0..255
        rgb = (int(r * 43 + 21), int(g * 43 + 21), int(b * 43 + 21))
        out.append({
            "hex": f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}",
            "weight": round(count / total, 4),
        })
    return out


# ── Pipeline ──────────────────────────────────────────────────────────────

class SlideIndexer:
    def __init__(self) -> None:
        self.s3 = S3Service()

    def index_pptx(
        self,
        *,
        user_id: str,
        source_file_id: str,
        s3_key: str,
        thumbnail_prefix: str,
    ) -> dict:
        """Full pipeline for one PPTX. Returns {slide_count, indexed} on success.

        thumbnail_prefix: S3 key prefix for storing slide PNGs, e.g. f"find/{user_id}/{source_file_id}".
        """
        log = logger.bind(source_file_id=source_file_id, user_id=user_id)
        log.info("indexing_started", s3_key=s3_key)

        with tempfile.TemporaryDirectory(prefix="slideforge-find-") as tmpdir:
            tmpdir_path = Path(tmpdir)
            pptx_path = tmpdir_path / "deck.pptx"
            pptx_path.write_bytes(self.s3.download_bytes(s3_key))

            # 1. Render slides
            png_dir = tmpdir_path / "pngs"
            pngs = render_pptx_to_pngs(pptx_path, png_dir)
            log.info("rendered", count=len(pngs))

            # 2. Extract text per slide
            slide_texts = extract_slide_texts(pptx_path)
            # Pad/truncate to match number of PNGs (PPTX text vs render count can drift on weird files)
            n = len(pngs)
            slide_texts = (slide_texts + [""] * n)[:n]

            # 3. OCR each PNG, compute embeddings, upload thumbnails, write to DB
            images: list[Image.Image] = []
            ocr_texts: list[str] = []
            color_palettes: list[list[dict]] = []
            thumb_keys: list[str] = []

            for idx, png_path in enumerate(pngs, start=1):
                slide_num = idx
                with Image.open(png_path) as img:
                    img.load()
                    img_rgb = img.convert("RGB")

                ocr_t = ocr_image(img_rgb)
                colors = dominant_colors(img_rgb)

                # Resize for thumbnail upload (web display) and CLIP input
                thumb = img_rgb.copy()
                thumb.thumbnail((1200, 1200), Image.LANCZOS)
                buf = io.BytesIO()
                thumb.save(buf, format="PNG", optimize=True)
                thumb_key = f"{thumbnail_prefix}/slide-{slide_num}.png"
                self.s3.upload_bytes(buf.getvalue(), thumb_key, content_type="image/png")

                images.append(img_rgb)
                ocr_texts.append(ocr_t)
                color_palettes.append(colors)
                thumb_keys.append(thumb_key)

            # Combined text for embedding: prioritize slide_text, fall back to OCR
            combined_for_embedding = [
                (slide_texts[i] + " " + ocr_texts[i]).strip() or "(blank slide)"
                for i in range(n)
            ]
            text_embeds = embed_text(combined_for_embedding)
            image_embeds = embed_image(images)

            log.info("embeddings_computed", n=n)

            # 4. Persist to slide_index. Wipe any stale rows for this source file first.
            with psycopg.connect(settings.database_url) as conn:
                register_vector(conn)
                with conn.cursor() as cur:
                    cur.execute(
                        'DELETE FROM slide_index WHERE "sourceFileId" = %s',
                        (source_file_id,),
                    )
                    for i in range(n):
                        cur.execute(
                            """
                            INSERT INTO slide_index
                                ("userId", "sourceFileId", "slideNumber", "thumbnailS3Key",
                                 "slideText", "ocrText", text_embedding, image_embedding, dominant_colors)
                            VALUES (%s, %s, %s, %s, %s, %s, %s, %s, %s::jsonb)
                            """,
                            (
                                user_id,
                                source_file_id,
                                i + 1,
                                thumb_keys[i],
                                slide_texts[i] or None,
                                ocr_texts[i] or None,
                                text_embeds[i].tolist(),
                                image_embeds[i].tolist(),
                                psycopg.types.json.Jsonb(color_palettes[i]),
                            ),
                        )
                conn.commit()

            log.info("indexing_complete", slides=n)
            return {"slide_count": n, "indexed": n}
