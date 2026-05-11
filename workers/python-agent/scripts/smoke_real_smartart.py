"""Smoke test for the real native SmartArt pipeline.

Builds a 1-slide deck with a Basic Cycle SmartArt populated from items. Saves
the deck via python-pptx, then runs the post-processor to inject the diagram
parts. Output: /tmp/smoke_real_smartart.pptx — open in PowerPoint Mac to
validate the SmartArt is treated as native (right-click → 'Convert to' /
'Change Colors' / 'Reset Graphic' should work).

Run:
    cd workers/python-agent && PYTHONPATH=. .venv/bin/python scripts/smoke_real_smartart.py
"""

from __future__ import annotations

import shutil
import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from app.preso_pro.smart_art import (
    inject_smart_arts,
    smart_art_cycle,
)
from app.preso_pro.planning.slide_spec import (
    CompositionRules,
    DeckContext,
    PaletteEntry,
    Typography,
    TypographyScale,
)


def make_ctx() -> DeckContext:
    palette = {
        "background":   PaletteEntry(hex="#FFFFFF", source="auto"),
        "surface":      PaletteEntry(hex="#F8FAFC", source="auto"),
        "primary":      PaletteEntry(hex="#5B8DEF", source="auto"),
        "accent_1":     PaletteEntry(hex="#F472B6", source="auto"),
        "accent_2":     PaletteEntry(hex="#34D399", source="auto"),
        "text_primary": PaletteEntry(hex="#0B1020", source="auto"),
        "text_muted":   PaletteEntry(hex="#64748B", source="auto"),
        "text_inverse": PaletteEntry(hex="#FFFFFF", source="auto"),
    }
    return DeckContext(
        palette=palette,
        typography=Typography(heading_font="Inter", body_font="Inter", scale=TypographyScale()),
        composition=CompositionRules(
            mood="vibrant-tech", background_mode="light",
            decoratives_allowed=[], decoratives_density="medium",
            min_negative_space=0.30,
        ),
        audience="marketing",
    )


def main() -> None:
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000

    blank_layout = prs.slide_layouts[6]  # blank layout in default master
    slide = prs.slides.add_slide(blank_layout)
    ctx = make_ctx()

    # Place a title above the SmartArt
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(9), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "Real native SmartArt — Basic Cycle"
    run.font.size = Pt(32)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0B, 0x10, 0x20)

    # Drop the SmartArt
    smart_art_cycle(
        slide, ctx,
        items=["Plan", "Build", "Measure", "Learn", "Iterate"],
        anchor="center",
        size_frac=0.62,
    )

    # Step 1: Save with python-pptx (this writes the marker rectangle but no diagram)
    with tempfile.TemporaryDirectory() as td:
        intermediate = Path(td) / "intermediate.pptx"
        prs.save(str(intermediate))

        # Step 2: Run post-processor to inject the SmartArt diagram parts
        pending = getattr(prs.part.package, "_pending_smart_art", [])
        print(f"Pending SmartArt requests: {len(pending)}")
        for sa in pending:
            print(f"  - layout={sa.layout_key} slide={sa.slide_index} marker={sa.marker_id}")

        out = Path("/tmp/smoke_real_smartart.pptx")
        inject_smart_arts(intermediate, pending, out)

    print(f"\nWrote {out}  ({out.stat().st_size:,} bytes)")
    print("\nOpen in PowerPoint to validate it renders as native SmartArt.")


if __name__ == "__main__":
    main()
