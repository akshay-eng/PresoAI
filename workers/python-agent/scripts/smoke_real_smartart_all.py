"""Smoke test: every supported visually-rich SmartArt layout, one per slide.

Output: /tmp/smoke_real_smartart_all.pptx
"""

from __future__ import annotations

import tempfile
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

from app.preso_pro.planning.slide_spec import (
    CompositionRules,
    DeckContext,
    PaletteEntry,
    Typography,
    TypographyScale,
)
from app.preso_pro.smart_art import (
    inject_smart_arts,
    smart_art_chevron_process,
    smart_art_chevron_rows,
    smart_art_column_blocks,
    smart_art_cycle,
    smart_art_hexagon_timeline,
    smart_art_horizontal_hierarchy,
    smart_art_org_chart,
    smart_art_trapezoid_blocks,
)


def make_ctx() -> DeckContext:
    palette = {
        "background":   PaletteEntry(hex="#FFFFFF", source="auto"),
        "surface":      PaletteEntry(hex="#F8FAFC", source="auto"),
        "primary":      PaletteEntry(hex="#5B8DEF", source="auto"),
        "accent_1":     PaletteEntry(hex="#F472B6", source="auto"),
        "accent_2":     PaletteEntry(hex="#34D399", source="auto"),
        "text_primary": PaletteEntry(hex="#0B1020", source="auto"),
        "text_muted":   PaletteEntry(hex="#64748B", source="auto"),
        "text_inverse": PaletteEntry(hex="#FFFFFF", source="auto"),
    }
    return DeckContext(
        palette=palette,
        typography=Typography(heading_font="Inter", body_font="Inter", scale=TypographyScale()),
        composition=CompositionRules(
            mood="vibrant-tech", background_mode="light",
            decoratives_allowed=[], decoratives_density="medium",
            min_negative_space=0.30,
        ),
        audience="marketing",
    )


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(12), Inches(0.85))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(28)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x0B, 0x10, 0x20)


def main() -> None:
    prs = Presentation()
    prs.slide_width = 12192000
    prs.slide_height = 6858000
    blank = prs.slide_layouts[6]
    ctx = make_ctx()

    # 1 — Continuous Cycle (cycle3) — 5 boxes around an arc
    s = prs.slides.add_slide(blank)
    add_title(s, "Continuous Cycle — How we iterate")
    smart_art_cycle(s, ctx,
        items=["Plan", "Build", "Measure", "Learn", "Iterate"],
        anchor="center", size_frac=0.72)

    # 2 — Organization Chart (orgChart1)
    s = prs.slides.add_slide(blank)
    add_title(s, "Organization Chart — Team structure")
    smart_art_org_chart(s, ctx,
        items=["CEO", "CTO", "CMO", "CFO"],
        anchor="center", size_frac=0.78)

    # 3 — Chevron rows (lProcess3) — 3 rows × 3 chevrons
    s = prs.slides.add_slide(blank)
    add_title(s, "Chevron Rows — Parallel workstreams")
    smart_art_chevron_rows(s, ctx,
        items=[
            "Plan", "Build", "Ship",        # row 1
            "Hire", "Train", "Scale",       # row 2
            "Discover", "Design", "Launch", # row 3
        ],
        anchor="center", size_frac=0.85)

    # 4 — Column blocks (lProcess2) — 3 cols × (title + 2 cells)
    s = prs.slides.add_slide(blank)
    add_title(s, "Column Blocks — Three pillars")
    smart_art_column_blocks(s, ctx,
        items=[
            "Speed",       "Sub-second responses", "Always-on global edge",
            "Trust",       "SOC2 compliant",       "Audit logs everywhere",
            "Insight",     "AI-driven anomalies",  "Auto-routed alerts",
        ],
        anchor="center", size_frac=0.78)

    # 5 — Horizontal hierarchy (hierarchy5) — 3-col tree
    s = prs.slides.add_slide(blank)
    add_title(s, "Hierarchy — Decision tree")
    smart_art_horizontal_hierarchy(s, ctx,
        items=[
            "Mission",       # root
            "Product",       # mid-tier 1
            "Operations",    # mid-tier 2
            "Engineering",   # leaves
            "Design",
            "Sales",
            "Support",
            "",  # padding
            "",
        ],
        anchor="center", size_frac=0.85)

    # 6 — Hexagon Timeline (HexagonTimeline) — 3 hexagons, title + caption each
    s = prs.slides.add_slide(blank)
    add_title(s, "Hexagon Timeline — Roadmap")
    smart_art_hexagon_timeline(s, ctx,
        items=[
            "Q1 2025", "Founded — closed seed round",
            "Q3 2025", "First 1,000 users on the platform",
            "Q2 2026", "Series A — enterprise tier launch",
        ],
        anchor="center", size_frac=0.85)

    # 7 — Accent Chevron Process (AccentHomeChevronProcess)
    s = prs.slides.add_slide(blank)
    add_title(s, "Chevron Process — Engagement flow")
    smart_art_chevron_process(s, ctx,
        items=[
            "Discover", "Workshops + interviews",
            "Build",    "Pilot deployment",
            "Scale",    "Org-wide rollout",
        ],
        anchor="center", size_frac=0.85)

    # 8 — Trapezoid blocks (hList6) — 3 columns
    s = prs.slides.add_slide(blank)
    add_title(s, "Trapezoid Blocks — Three core values")
    smart_art_trapezoid_blocks(s, ctx,
        items=[
            "Customer first",  "We obsess over customer outcomes",  "Every decision starts here",
            "Speed wins",      "We ship and iterate weekly",        "Speed beats polish in early markets",
            "Trust always",    "Security and compliance from day one", "We earn trust by being predictable",
        ],
        anchor="center", size_frac=0.72)

    # Save + post-process
    with tempfile.TemporaryDirectory() as td:
        intermediate = Path(td) / "intermediate.pptx"
        prs.save(str(intermediate))
        pending = getattr(prs.part.package, "_pending_smart_art", [])
        print(f"Pending SmartArt requests: {len(pending)}")
        for sa in pending:
            print(f"  slide {sa.slide_index}: {sa.layout_key}")

        out = Path("/tmp/smoke_real_smartart_all.pptx")
        inject_smart_arts(intermediate, pending, out)

    print(f"\nWrote {out}  ({out.stat().st_size:,} bytes)")


if __name__ == "__main__":
    main()
